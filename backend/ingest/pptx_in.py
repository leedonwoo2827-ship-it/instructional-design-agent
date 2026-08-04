# -*- coding: utf-8 -*-
"""PPTX → 슬라이드플랜 추출.

instructional-design-agent(`core/deck_builder.py`)가 뽑는 덱을 주 대상으로 하되,
임의의 덱에서도 깨지지 않게 방어적으로 읽는다.

관찰된 덱 구조 (교육공학 2주차, 102장):
  - 텍스트 런 2개              → 섹션 구분 슬라이드
  - 제목 + 요약문 + `▸` 반복   → 불릿형
  - 제목 + 요약문 + 1/레이블/설명 반복 → 번호 카드형
  - 제목 + 요약문 + 좌우 레이블/설명   → 2단 비교형

두 번째 텍스트가 내레이션형 한 문장인 경우가 있어 `spine`으로 분리한다.
단 2026-07-28 첨삭 이후 덱은 명사형 어미로 바뀌어 spine이 없을 수 있으므로 **optional**이다.
"""
from __future__ import annotations

import re
from dataclasses import dataclass, field, asdict
from pathlib import Path

BULLET_MARK = "▸"
# "1", "2" … 처럼 카드 번호로만 이루어진 런
_NUM_ONLY = re.compile(r"^\d{1,2}$")
# 서술형 종결 — spine 후보 판별용
_SENTENCE_END = re.compile(r"(다|요|까|죠)[.!?]?$")


@dataclass
class Slide:
    index: int                       # 1-based
    kind: str                        # section | bullets | cards | compare | content
    title: str = ""
    spine: str = ""                  # 있으면 내레이션 뼈대 (없을 수 있음)
    bullets: list[str] = field(default_factory=list)
    cards: list[dict] = field(default_factory=list)   # {"label","desc"}
    has_image: bool = False
    has_table: bool = False

    def to_dict(self) -> dict:
        return asdict(self)

    def text_for_prompt(self) -> str:
        """LLM 프롬프트에 넣을 슬라이드 요약."""
        parts = [f"제목: {self.title}"]
        if self.spine:
            parts.append(f"핵심: {self.spine}")
        for b in self.bullets:
            parts.append(f"- {b}")
        for c in self.cards:
            lbl, desc = c.get("label", ""), c.get("desc", "")
            parts.append(f"- {lbl}: {desc}" if desc else f"- {lbl}")
        if self.has_image:
            parts.append("(이미지 있음)")
        return "\n".join(parts)


def _runs(shape) -> list[str]:
    """도형에서 텍스트를 문단 단위로 뽑는다 (런 분할은 서식 때문이라 무의미)."""
    out: list[str] = []
    if not getattr(shape, "has_text_frame", False):
        return out
    for para in shape.text_frame.paragraphs:
        t = "".join(r.text for r in para.runs).strip()
        if t:
            out.append(t)
    return out


def _walk(shapes) -> tuple[list[str], bool, bool]:
    """그룹 도형까지 재귀. (텍스트, 이미지있음, 표있음)"""
    texts: list[str] = []
    has_img = has_tbl = False
    for sh in shapes:
        st = sh.shape_type
        if st is not None and int(st) == 6:          # GROUP
            t, i, b = _walk(sh.shapes)
            texts += t
            has_img |= i
            has_tbl |= b
            continue
        if getattr(sh, "has_table", False):
            has_tbl = True
            for row in sh.table.rows:
                for cell in row.cells:
                    t = cell.text.strip()
                    if t:
                        texts.append(t)
            continue
        if st is not None and int(st) == 13:         # PICTURE
            has_img = True
            continue
        texts += _runs(sh)
    return texts, has_img, has_tbl


def _classify(texts: list[str]) -> tuple[str, str, str, list[str], list[dict]]:
    """텍스트 목록 → (kind, title, spine, bullets, cards)"""
    title = texts[0] if texts else ""
    rest = texts[1:]

    if not rest:
        return "section", title, "", [], []

    # 두 번째 텍스트가 서술형 문장이면 spine 후보
    spine = ""
    if _SENTENCE_END.search(rest[0]) and len(rest[0]) > 8:
        spine = rest[0]
        rest = rest[1:]

    if not rest:
        return "section", title, spine, [], []

    # 번호 카드형: "1" 다음에 레이블, 그 다음 설명
    if any(_NUM_ONLY.match(t) for t in rest):
        cards: list[dict] = []
        i = 0
        while i < len(rest):
            if _NUM_ONLY.match(rest[i]):
                label = rest[i + 1] if i + 1 < len(rest) else ""
                desc = rest[i + 2] if i + 2 < len(rest) and not _NUM_ONLY.match(rest[i + 2]) else ""
                cards.append({"label": label, "desc": desc})
                i += 3 if desc else 2
            else:
                i += 1
        if cards:
            return "cards", title, spine, [], cards

    # 불릿형: ▸ 가 별도 런으로 오거나("▸", "내용") 문단 앞에 붙어 온다("▸ 내용").
    if any(t == BULLET_MARK or t.startswith(BULLET_MARK) for t in rest):
        bullets = [t.lstrip(BULLET_MARK).strip() for t in rest if t != BULLET_MARK]
        bullets = [b for b in bullets if b]
        return "bullets", title, spine, bullets, []

    # 2단 비교형: 짧은 레이블 + 긴 설명이 쌍으로 반복될 때만.
    # (길이 조건이 없으면 평범한 문단 나열까지 전부 compare로 삼켜버린다)
    if len(rest) >= 4 and len(rest) % 2 == 0:
        pairs = [(rest[i], rest[i + 1]) for i in range(0, len(rest), 2)]
        if all(len(lbl) <= 30 and len(desc) > len(lbl) for lbl, desc in pairs):
            return "compare", title, spine, [], [{"label": l, "desc": d} for l, d in pairs]

    return "content", title, spine, rest, []


def extract(pptx_path: str | Path) -> list[Slide]:
    from pptx import Presentation

    prs = Presentation(str(pptx_path))
    slides: list[Slide] = []
    for i, s in enumerate(prs.slides, start=1):
        texts, has_img, has_tbl = _walk(s.shapes)
        # 노트가 있으면 그게 최고 품질의 spine이다 (deck_builder는 안 넣지만 외부 덱은 넣을 수 있음)
        note = ""
        if s.has_notes_slide:
            note = (s.notes_slide.notes_text_frame.text or "").strip()

        kind, title, spine, bullets, cards = _classify(texts)
        slides.append(Slide(
            index=i, kind=kind, title=title,
            spine=note or spine,
            bullets=bullets, cards=cards,
            has_image=has_img, has_table=has_tbl,
        ))
    return slides


# ── 슬라이드플랜(ID agent STEP 4) → Slide ────────────────────────────────
# instructional-design-agent `core/deck_builder.py` 가 build_deck() 에 넣는 그 구조다.
# PPTX를 다시 파싱하는 것보다 정확하다 — 분류 추측이 필요 없다.
#
#   {"type":"section|photo|process|cards|compare|table|bullets|cover",
#    "title":..., "chip":"핵심 메시지 한 문장(명사형)",
#    "bullets":[...], "items":[{"label","desc"|"lines"}],
#    "headers":[...], "rows":[[...]], "image_query":..., "numbered":bool}
#
# ★ `chip` 은 명사형이라 그대로 TTS 하면 안 된다. spine 으로만 쓰고 대본 생성을 반드시 거친다.
_PLAN_KIND = {
    "cover": "cover", "section": "section", "photo": "photo",
    "process": "process", "cards": "cards", "compare": "compare",
    "table": "table", "bullets": "bullets",
}


def from_plan(plan: list[dict]) -> list[Slide]:
    slides: list[Slide] = []
    for i, sp in enumerate(plan, start=1):
        kind = _PLAN_KIND.get(str(sp.get("type", "")).strip(), "bullets")

        cards: list[dict] = []
        for it in (sp.get("items") or []):
            if not isinstance(it, dict):
                cards.append({"label": str(it), "desc": ""})
                continue
            desc = it.get("desc") or ""
            if not desc and it.get("lines"):          # compare 형은 lines 로 온다
                desc = " ".join(str(x) for x in it["lines"])
            cards.append({"label": str(it.get("label", "")), "desc": str(desc)})

        bullets = [str(b) for b in (sp.get("bullets") or []) if str(b).strip()]

        # 표는 헤더+행을 카드처럼 눌러 담는다 — 대본은 패턴을 설명하게 프롬프트가 지시한다
        if kind == "table" and sp.get("rows"):
            headers = [str(h) for h in (sp.get("headers") or [])]
            for row in sp["rows"]:
                cells = [str(c) for c in row]
                label = cells[0] if cells else ""
                rest = cells[1:]
                desc = ", ".join(
                    f"{headers[j+1]}: {v}" if j + 1 < len(headers) else v
                    for j, v in enumerate(rest))
                cards.append({"label": label, "desc": desc})

        slides.append(Slide(
            index=i, kind=kind,
            title=str(sp.get("title", "")),
            spine=str(sp.get("chip") or sp.get("subtitle") or ""),
            bullets=bullets, cards=cards,
            has_image=bool(sp.get("image_query")),
            has_table=bool(sp.get("rows")),
        ))
    return slides


def load_plan(path: str | Path) -> list[Slide]:
    import json
    data = json.loads(Path(path).read_text(encoding="utf-8"))
    plan = data.get("slides") if isinstance(data, dict) else data
    return from_plan(plan or [])


if __name__ == "__main__":
    import sys, json
    if sys.argv[1].lower().endswith(".json"):
        ss = load_plan(sys.argv[1])
    else:
        ss = extract(sys.argv[1])
    from collections import Counter
    print(f"슬라이드 {len(ss)}장  유형: {dict(Counter(s.kind for s in ss))}")
    print(f"spine 보유: {sum(1 for s in ss if s.spine)}장 / 이미지: {sum(1 for s in ss if s.has_image)}장")
    for s in ss[:4]:
        print("\n" + "-" * 60)
        print(s.text_for_prompt())
