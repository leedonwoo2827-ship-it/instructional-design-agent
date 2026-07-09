# -*- coding: utf-8 -*-
"""PPT 개요(마크다운) → 실제 .pptx 변환. (회사 양식 템플릿 상속 지원)

Anthropic pptx 스킬의 '템플릿 기반 생성' 방식을 적용:
- 회사 양식 .pptx 가 있으면 그것을 베이스로 열어(테마·마스터·폰트·레이아웃 상속)
  기존 예시 슬라이드는 비우고, 양식의 레이아웃으로 새 슬라이드를 채운다.
- 양식이 없으면 기본 16:9 프레젠테이션으로 생성한다.

SYS_SCRIPT_PPT 개요 형식(### 슬라이드 N — 제목 + 불릿 + 발표자 노트)을 파싱한다.
"""
from __future__ import annotations

import io
import os
import re
from typing import Optional


def _clean(text: str) -> str:
    text = re.sub(r"\*\*|__|`|~~", "", text)
    text = re.sub(r"^[\-\*\+•·]\s*", "", text)
    return text.strip()


def _remove_all_slides(prs) -> None:
    """템플릿에 포함된 예시 슬라이드 제거(마스터·레이아웃은 보존).

    sldId 와 그 관계(rId)를 함께 끊는다. 도달 불가가 된 슬라이드 파트는
    저장 시 직렬화되지 않으므로 partname 충돌(Duplicate name)이 발생하지 않는다.
    """
    from pptx.oxml.ns import qn

    part = prs.part
    id_lst = prs.slides._sldIdLst
    for sid in list(id_lst):
        rid = sid.get(qn("r:id"))
        if rid and rid in part.rels:
            try:
                part.drop_rel(rid)
            except Exception:  # noqa: BLE001
                pass
        id_lst.remove(sid)


def _find_layout(prs, keywords, fallback_idx):
    for layout in prs.slide_layouts:
        nm = (layout.name or "").lower()
        if any(k in nm for k in keywords):
            return layout
    try:
        return prs.slide_layouts[fallback_idx]
    except Exception:  # noqa: BLE001
        return prs.slide_layouts[0]


def _body_placeholder(slide):
    """제목(idx 0) 이외의 첫 텍스트 플레이스홀더."""
    for ph in slide.placeholders:
        try:
            if ph.placeholder_format.idx != 0 and ph.has_text_frame:
                return ph
        except Exception:  # noqa: BLE001
            continue
    return None


def _content_phs(slide):
    """제목(idx 0) 이외의 텍스트 플레이스홀더들을 idx 순으로 반환(2단 배치용)."""
    out = []
    for ph in slide.placeholders:
        try:
            if ph.placeholder_format.idx != 0 and ph.has_text_frame:
                out.append(ph)
        except Exception:  # noqa: BLE001
            continue
    return out


def _fill(ph, lines):
    """플레이스홀더 텍스트프레임에 여러 줄 채우기(첫 줄 굵게)."""
    tf = ph.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, line in enumerate(lines[:10]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line[:220]
        if i == 0:
            p.font.bold = True


def outline_to_pptx(md: str, deck_title: str = "강의 슬라이드",
                    template_path: Optional[str] = None) -> Optional[bytes]:
    """개요 마크다운 → .pptx 바이트. python-pptx 미설치 시 None."""
    try:
        from pptx import Presentation
        from pptx.util import Pt
    except Exception:  # noqa: BLE001
        return None

    use_template = bool(template_path and os.path.isfile(template_path))
    if use_template:
        try:
            prs = Presentation(template_path)
            _remove_all_slides(prs)
        except Exception:  # noqa: BLE001
            prs = Presentation()
            use_template = False
    else:
        prs = Presentation()

    if not use_template:
        prs.slide_width = Pt(960)
        prs.slide_height = Pt(540)

    title_layout = _find_layout(prs, ["title slide", "제목 슬라이드", "표지", "cover"], 0)
    body_layout = _find_layout(prs, ["title and content", "제목 및 내용", "content", "내용", "본문"], 1)
    section_layout = _find_layout(prs, ["section header", "구역 머리글", "섹션", "구역"], 2)
    two_layout = _find_layout(prs, ["콘텐츠 2개", "두 개의 콘텐츠", "two content", "comparison", "비교", "2개"], 3)

    # ── 표지 ──
    s = prs.slides.add_slide(title_layout)
    if s.shapes.title is not None:
        s.shapes.title.text = deck_title
    sub = _body_placeholder(s)
    if sub is not None:
        sub.text_frame.text = "교수설계 가이드 에이전트 · Mayer 멀티미디어 원리 기반"

    # ── 슬라이드 블록 파싱 ──
    label_re = re.compile(r"^\s*[-*+]?\s*\*{0,2}([^*:：]{1,24})\*{0,2}\s*[:：]\s*(.*)$")
    blocks = re.split(r"(?m)^\s*#{2,3}\s+", md or "")
    slide_blocks = [b for b in blocks if re.match(r"\s*슬라이드", b)]
    if not slide_blocks:
        slide_blocks = [b for b in blocks if b.strip()][:20]

    for block in slide_blocks:
        lines = block.splitlines()
        header = lines[0].strip() if lines else "슬라이드"
        title = re.sub(r"^슬라이드\s*\d+\s*[—\-:：]\s*", "", header).strip() or header

        layout_hint, body, notes, mode = "", [], [], "body"
        for ln in lines[1:]:
            t = ln.strip()
            if not t:
                continue
            m = label_re.match(t)
            if m:
                label, rest = m.group(1).strip(), m.group(2).strip()
                if "레이아웃" in label:
                    layout_hint = rest
                    mode = "body"
                    continue
                if "발표자" in label and "노트" in label:
                    mode = "notes"
                    if rest:
                        notes.append(_clean(rest))
                    continue
                if "시각자료" in label or "강조" in label:
                    if rest:
                        notes.append(f"({label}) {_clean(rest)}")
                    mode = "body"
                    continue
                if "핵심" in label and "메시지" in label:
                    if rest:
                        body.append(_clean(rest))
                    mode = "body"
                    continue
                if "본문" in label:
                    # "불릿 3~5개" 같은 안내 문구는 건너뜀
                    if rest and not re.search(r"불릿|\d\s*~?\s*\d*\s*개", rest):
                        body.append(_clean(rest))
                    mode = "body"
                    continue
                # 기타 라벨: 값만 본문/노트에
                if rest:
                    (notes if mode == "notes" else body).append(_clean(rest))
                continue
            (notes if mode == "notes" else body).append(_clean(t))

        hint = layout_hint
        is_section = bool(hint) and ("섹션" in hint or "표지" in hint)
        is_two = any(k in hint for k in ["2단", "두 단", "2 단", "two", "2x2", "2X2", "그리드", "grid", "비교", "콘텐츠 2"])

        if is_section:
            layout = section_layout
        elif is_two:
            layout = two_layout
        else:
            layout = body_layout
        slide = prs.slides.add_slide(layout)
        if slide.shapes.title is not None:
            slide.shapes.title.text = title[:120]

        cps = _content_phs(slide)
        if is_two and len(cps) >= 2 and len(body) >= 2:
            # 본문을 좌/우 두 칸에 나눠 배치
            mid = (len(body) + 1) // 2
            _fill(cps[0], body[:mid])
            _fill(cps[1], body[mid:])
        elif cps and body:
            _fill(cps[0], body)
        # 발표자 노트는 넣지 않는다(요청).

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
