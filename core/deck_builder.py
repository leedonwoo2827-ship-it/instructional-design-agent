# -*- coding: utf-8 -*-
"""디자인 슬라이드 빌더 — 슬라이드플랜(JSON) → 디자인된 .pptx 바이트.

네이비(#1E2761) + 앰버(#F2A900) 디자인 시스템. 회사 템플릿을 베이스로 열어
(테마·마스터·로고 상속) 예시 슬라이드를 비우고, 모든 슬라이드를 Blank 레이아웃에
도형으로 직접 배치한다.

타입: cover / section / photo / bullets / process / cards / compare / table
      / quiz / objectives / agenda / closing / stat
python-pptx 미설치 시 build_deck 은 None 을 돌려준다(그레이스풀).

디자인 규칙(docs/디자인-토큰.md 참고)
  · 색은 네이비 단일 계조 + 앰버 1점. 무지개 로테이션 금지.
  · 앰버는 채움·마커 전용. 흰 배경 위 본문 텍스트 색으로 쓰지 않는다(대비 2.0:1).
  · 그림자는 사진·카드·비교 컬럼에만. 칩·바탕 도형은 평면.
  · 모든 본문은 _fit() 을 거쳐 박스를 넘지 않는다.
"""
from __future__ import annotations

import io
import json
import math
import os
import re
from pathlib import Path
from typing import Callable, Dict, List, Optional

from core import pptx_font

# ── 캔버스 ────────────────────────────────────────────────────────────────
SLIDE_W, SLIDE_H = 13.33, 7.5
MARGIN = 0.6
CONTENT_W = 12.13                      # 전폭 콘텐츠
CONTENT_R = MARGIN + CONTENT_W         # 우측 끝 12.73

ASSETS_DIR = Path(__file__).resolve().parent.parent / "assets"

# ── 색 — 네이비 단일 계조 + 앰버 ───────────────────────────────────────────
_HEX = {
    "navy":     (0x1E, 0x27, 0x61),   # 제목·주요 채움 (흰 글씨 OK)
    "navy_75":  (0x56, 0x5D, 0x88),   # 2단계
    "navy_55":  (0x83, 0x88, 0xA8),   # 3단계·키커·캡션
    "navy_30":  (0xBB, 0xBE, 0xD0),   # 비활성·대형 숫자
    "navy_12":  (0xE4, 0xE5, 0xEC),   # 헤어라인 링
    "chip":     (0xEE, 0xF2, 0xFB),   # 존 배경·표 짝수행·사진 자리 워시
    "amber":    (0xF2, 0xA9, 0x00),   # 마커·강조 채움 (텍스트 금지)
    "amber_dk": (0xB5, 0x7F, 0x00),   # 앰버 계열 텍스트가 꼭 필요할 때
    "amber_wash": (0xFD, 0xF3, 0xDC),  # 강조 존 배경
    "ink":      (0x2B, 0x34, 0x40),   # 본문 글자
    "navy_dk":  (0x1E, 0x27, 0x61),   # 글자가 얹히는 밴드 배경 (기본은 navy 와 동일)
    "on_navy":  (0xA9, 0xB0, 0xCE),   # 네이비 위 보조 글자
    "white":    (0xFF, 0xFF, 0xFF),
}

# 순서 있는 것(process)·병렬(cards) 모두 단일 계조. 무지개 금지.
_RAMP = ("navy", "navy_75", "navy_55", "navy_30")

# ── 타이포 스케일 ─────────────────────────────────────────────────────────
#   size(pt) / weight(r|sb|b) / color / spc(em 자간) / line(행간)
TYPE = {
    "kicker":      dict(size=11.0, weight="sb", color="navy_55", spc=0.14),
    "kicker_on":   dict(size=11.0, weight="sb", color="on_navy", spc=0.14),
    "cover_title": dict(size=40.0, weight="b",  color="navy", spc=-0.025, line=1.14),
    "cover_sub":   dict(size=17.0, weight="r",  color="navy_55", line=1.35),
    "sec_title":   dict(size=32.0, weight="b",  color="navy", spc=-0.02, line=1.18),
    "title":       dict(size=28.0, weight="b",  color="navy", spc=-0.02, line=1.16),
    "chip":        dict(size=15.0, weight="sb", color="navy", line=1.25),
    "bullet":      dict(size=15.0, weight="r",  color="ink", line=1.42),
    "bullet2":     dict(size=13.0, weight="r",  color="navy_55", line=1.40),
    "card_label":  dict(size=15.5, weight="sb", color="navy", line=1.20),
    "card_desc":   dict(size=12.5, weight="r",  color="ink", line=1.38),
    "node_label":  dict(size=14.5, weight="sb", color="white", line=1.18),
    "node_desc":   dict(size=11.0, weight="r",  color="ink", line=1.30),
    "th":          dict(size=12.0, weight="sb", color="white"),
    "td":          dict(size=11.5, weight="r",  color="ink", line=1.30),
    "caption":     dict(size=9.5,  weight="r",  color="navy_55"),
    "footer":      dict(size=9.5,  weight="r",  color="navy_55"),
    "pageno":      dict(size=10.0, weight="sb", color="navy_30"),
    "big_num":     dict(size=34.0, weight="b",  color="navy_30", spc=-0.02),
    "sec_num":     dict(size=96.0, weight="b",  color="navy_12", spc=-0.03),
    "stat_num":    dict(size=52.0, weight="b",  color="navy", spc=-0.03),
    "answer":      dict(size=13.0, weight="sb", color="navy"),
}

# ── 세로 그리드 ───────────────────────────────────────────────────────────
KICKER_Y = 0.52
TITLE_Y, TITLE_H = 0.80, 0.94
RULE_Y, RULE_W, RULE_H = 1.66, 0.62, 0.05
CHIP_Y, CHIP_H = 1.88, 0.86
BODY_TOP = 2.95          # 콘텐츠 상단 기준선
BODY_BOTTOM = 6.62       # 콘텐츠 하단 한계
FOOT_LINE_Y = 6.90       # 러닝 헤어라인
FOOT_TEXT_Y = 6.98

# 표준 시각 반경(인치)
R_CARD = 0.10
R_PHOTO = 0.10

# 현재 빌드의 폰트 세트(build_deck 진입 시 설정)
_FS = pptx_font.SYSTEM_SET


def _log(msg: str) -> None:
    """콘솔 코드페이지(cp949 등)에 없는 글자로 로그가 죽지 않게 한다.

    렌더 함수 안에서 print 가 UnicodeEncodeError 를 던지면 build_deck 의
    try/except 가 이를 렌더 실패로 오인해 폴백시킨다 — 실제로 겪은 문제.
    """
    try:
        print(msg, flush=True)
    except Exception:  # noqa: BLE001
        try:
            enc = getattr(getattr(__import__("sys"), "stdout", None), "encoding",
                          None) or "ascii"
            print(msg.encode(enc, "replace").decode(enc, "replace"), flush=True)
        except Exception:  # noqa: BLE001
            pass


def _rgb(name):
    from pptx.dml.color import RGBColor
    return RGBColor(*_HEX[name])


def _hexstr(name: str) -> str:
    return "%02X%02X%02X" % _HEX[name]


def luminance(name: str) -> float:
    """상대 휘도(0~1). WCAG 식 — 감마 보정을 해야 사람 눈과 맞는다."""
    def lin(c: float) -> float:
        c /= 255.0
        return c / 12.92 if c <= 0.03928 else ((c + 0.055) / 1.055) ** 2.4
    r, g, b = _HEX.get(name, (0, 0, 0))
    return 0.2126 * lin(r) + 0.7152 * lin(g) + 0.0722 * lin(b)


def band(name: str, min_ratio: float = 4.5) -> str:
    """글자를 얹을 밴드의 채움색 — 대비가 모자라면 `<name>_dk` 로 바꾼다.

    템플릿마다 주색의 밝기가 다르다. 기본 배색의 navy(#1E2761)는 흰글씨와
    13.8:1 이라 그대로 써도 되지만, 리디자인의 navy(#2E93D9)는 3.3:1 이고
    amber(#4CB782)는 2.5:1 로 어떤 글자색을 얹어도 안 읽힌다.

    그래서 **모자랄 때만** 진한 짝을 쓴다. 넉넉한 템플릿은 손대지 않으므로
    기존 덱의 외형이 바뀌지 않는다.
    """
    def _best(bg: str) -> float:
        lb = luminance(bg)
        def r(a, b):
            hi, lo = max(a, b), min(a, b)
            return (hi + 0.05) / (lo + 0.05)
        return max(r(luminance("white"), lb), r(luminance("ink"), lb))

    if _best(name) >= min_ratio:
        return name
    dk = f"{name}_dk"
    if dk in _HEX and _best(dk) > _best(name):
        return dk
    return name


def on_white(name: str, min_ratio: float = 4.5) -> str:
    """**흰 바탕 위 글자색** — 대비가 모자라면 `<name>_dk`, 그래도 안 되면 `ink`.

    band() 와 헷갈리면 안 된다. band() 는 *배경* 을 고르고, 이건 *글자* 를 고른다.
    칩이 배경(연한 알약)을 잃고 흰 바탕에 놓이면서 필요해졌다 —
    리디자인의 navy(#2E93D9)는 연한 칩 위에서는 괜찮았지만
    흰 바탕에서는 3.34:1 로 떨어진다.
    """
    def ratio(a: float, b: float) -> float:
        hi, lo = max(a, b), min(a, b)
        return (hi + 0.05) / (lo + 0.05)

    w = luminance("white")
    for cand in (name, f"{name}_dk", "ink"):
        try:
            if ratio(luminance(cand), w) >= min_ratio:
                return cand
        except Exception:  # noqa: BLE001 — 템플릿에 없는 키
            continue
    return "ink"


def on_color(bg: str, light: str = "white", dark: str = "ink") -> str:
    """배경 위에 얹을 글자색 — 명암비가 큰 쪽을 고른다.

    색을 박아 두면 템플릿에서 배경색을 바꿀 때 대비가 무너진다. 실제로 앰버
    밴드에는 어두운 글씨(6.9:1)가 맞고 흰 글씨(2.1:1)는 안 읽히는데, 초록 밴드는
    그 반대다. 그래서 고정값 대신 계산한다.

    ★ 어두운 쪽 기본값은 `navy` 가 아니라 **`ink`(본문 글자색)** 다. `navy` 는
      템플릿에 따라 진한 남색일 수도, 중간톤 블루일 수도 있다 — 리디자인 배색에서
      navy=#2E93D9 였고, 그걸 글자색으로 쓰면 연한 배경에서 2.8:1 로 떨어졌다.
      `ink` 는 어느 템플릿에서든 '가장 읽히는 글자색' 이라는 뜻이 유지된다.
    """
    def ratio(a: float, b: float) -> float:
        hi, lo = max(a, b), min(a, b)
        return (hi + 0.05) / (lo + 0.05)
    lb = luminance(bg)
    return light if ratio(luminance(light), lb) >= ratio(luminance(dark), lb) else dark


# ── 저수준 헬퍼 ───────────────────────────────────────────────────────────
def _remove_all_slides(prs) -> None:
    from pptx.oxml.ns import qn
    part = prs.part
    id_lst = prs.slides._sldIdLst
    for sid in list(id_lst):
        rid = sid.get(qn("r:id"))
        if rid and rid in part.rels:
            try:
                part.drop_rel(rid)
            except Exception:  # noqa: BLE001
                pass
        id_lst.remove(sid)


def _blank_layout(prs):
    for lay in prs.slide_layouts:
        nm = (lay.name or "").lower()
        if "blank" in nm or "빈" in (lay.name or ""):
            return lay
    return min(prs.slide_layouts, key=lambda l: len(l.placeholders))


def _strip_placeholders(slide) -> None:
    """레이아웃에서 상속된 빈 플레이스홀더 제거(마스터 로고·장식은 유지)."""
    for ph in list(slide.placeholders):
        try:
            ph._element.getparent().remove(ph._element)
        except Exception:  # noqa: BLE001
            pass


def _no_deco(shape):
    """도형 기본 그림자·테두리 제거."""
    try:
        shape.line.fill.background()
    except Exception:  # noqa: BLE001
        pass
    try:
        shape.shadow.inherit = False
    except Exception:  # noqa: BLE001
        pass


def _adj(radius_in: float, w: float, h: float) -> int:
    """시각 반경(inch)을 roundRect adj 값으로. 도형 크기가 달라도 라운드가 같아진다."""
    base = max(min(w, h), 0.01)
    return int(max(0.0, min(radius_in / base, 0.5)) * 100000)


def _set_adj(shape, val: int) -> None:
    from pptx.oxml.ns import qn
    geom = shape._element.spPr.find(qn("a:prstGeom"))
    if geom is None:
        return
    av = geom.find(qn("a:avLst"))
    if av is None:
        av = geom.makeelement(qn("a:avLst"), {})
        geom.append(av)
    for gd in list(av):
        av.remove(gd)
    gd = av.makeelement(qn("a:gd"), {"name": "adj", "fmla": f"val {val}"})
    av.append(gd)


def _shadow(shape, *, blur=0.14, dist=0.035, alpha=12, color="navy") -> None:
    """부드러운 아래 방향 그림자. 사진·카드·비교 컬럼에만 쓴다."""
    from pptx.oxml.ns import qn
    try:
        spPr = shape._element.spPr
        for el in spPr.findall(qn("a:effectLst")):
            spPr.remove(el)
        eff = spPr.makeelement(qn("a:effectLst"), {})
        shd = eff.makeelement(qn("a:outerShdw"), {
            "blurRad": str(int(blur * 914400)),
            "dist": str(int(dist * 914400)),
            "dir": "5400000", "rotWithShape": "0"})
        clr = shd.makeelement(qn("a:srgbClr"), {"val": _hexstr(color)})
        clr.append(clr.makeelement(qn("a:alpha"), {"val": str(int(alpha * 1000))}))
        shd.append(clr)
        eff.append(shd)
        spPr.append(eff)
    except Exception:  # noqa: BLE001
        pass


def _rect(slide, x, y, w, h, fill):
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    _no_deco(shp)
    shp.fill.solid()
    shp.fill.fore_color.rgb = _rgb(fill)
    return shp


def _rrect(slide, x, y, w, h, fill=None, radius=R_CARD, line=None, line_w=None,
           pill=False):
    """라운드 사각형. radius 는 '인치 단위 시각 반경'(비율 아님). pill=True 면 알약."""
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches, Pt
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    _no_deco(shp)
    if fill is not None:
        shp.fill.solid()
        shp.fill.fore_color.rgb = _rgb(fill)
    else:
        shp.fill.background()
    if line is not None:
        shp.line.color.rgb = _rgb(line)
        shp.line.width = Pt(line_w or 1.0)
    _set_adj(shp, 50000 if pill else _adj(radius, w, h))
    return shp


def _oval(slide, x, y, d, fill):
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y),
                                 Inches(d), Inches(d))
    _no_deco(shp)
    shp.fill.solid()
    shp.fill.fore_color.rgb = _rgb(fill)
    return shp


def _arrow(slide, x, y, w, h, fill="amber"):
    """단계·순서 연결자."""
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches
    shp = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y),
                                 Inches(w), Inches(h))
    _no_deco(shp)
    shp.fill.solid()
    shp.fill.fore_color.rgb = _rgb(fill)
    return shp


def _text(slide, x, y, w, h, anchor="t"):
    from pptx.util import Inches
    from pptx.enum.text import MSO_ANCHOR
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = {"t": MSO_ANCHOR.TOP, "m": MSO_ANCHOR.MIDDLE,
                          "b": MSO_ANCHOR.BOTTOM}.get(anchor, MSO_ANCHOR.TOP)
    return tf


def _para(tf, first, align=None, space_after=4, line=1.05):
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    if align:
        p.alignment = {"c": PP_ALIGN.CENTER, "l": PP_ALIGN.LEFT,
                       "r": PP_ALIGN.RIGHT}[align]
    p.space_after = Pt(space_after)
    p.space_before = Pt(0)
    try:
        p.line_spacing = line
    except Exception:  # noqa: BLE001
        pass
    return p


def _run(p, text, size, color, weight="r", spc=0.0):
    from pptx.util import Pt
    r = p.add_run()
    r.text = str(text)
    r.font.size = Pt(size)
    r.font.bold = _FS.is_bold(weight)
    r.font.color.rgb = _rgb(color)
    pptx_font.apply_face(r, _FS.face(weight), spc_em=spc, size_pt=size)
    return r


def _T(p, text, role, *, size=None, color=None, weight=None, spc=None):
    """TYPE 롤로 런 추가. 개별 인자로 덮어쓸 수 있다."""
    t = TYPE[role]
    return _run(p, text,
                t["size"] if size is None else size,
                t["color"] if color is None else color,
                t["weight"] if weight is None else weight,
                t.get("spc", 0.0) if spc is None else spc)


def _hang(p, marL_in=0.24):
    """내어쓰기 — ▸ 는 왼쪽으로 튀고 줄바꿈 글자는 텍스트 아래로 정렬."""
    emu = str(int(marL_in * 914400))
    pPr = p._p.get_or_add_pPr()
    pPr.set("marL", emu)
    pPr.set("indent", "-" + emu)


# ── 텍스트 오버플로 보호 ───────────────────────────────────────────────────
def _measure(text: str, pt: float) -> float:
    """문자열의 근사 렌더 폭(inch). 한글·전각 1.0em, 라틴 0.54em, 공백 0.28em."""
    em = pt / 72.0
    w = 0.0
    for ch in str(text):
        o = ord(ch)
        if ch == " ":
            w += 0.28
        elif (0x1100 <= o <= 0x11FF or 0x2E80 <= o <= 0xA4CF
              or 0xAC00 <= o <= 0xD7A3 or 0xF900 <= o <= 0xFAFF
              or 0xFF00 <= o <= 0xFF60):
            w += 1.0
        elif ch.isdigit() or ch.isupper():
            w += 0.58
        else:
            w += 0.52
    return w * em


def _needed_h(lines: List[str], box_w: float, pt: float, line: float,
              space_after: float) -> float:
    """문단들이 차지할 높이(inch) 추정."""
    usable = max(box_w - 0.12, 0.4)
    rows = 0
    for t in lines:
        rows += max(1, math.ceil(_measure(t, pt) / usable - 1e-6))
    return rows * (pt / 72.0) * line + len(lines) * (space_after / 72.0)


def _fit(lines: List[str], box_w: float, box_h: float, base_pt: float,
         line: float = 1.4, space_after: float = 0.0,
         min_pt: float = 10.5, label: str = "") -> float:
    """박스에 들어가는 최대 폰트 크기(0.5pt 단위). 최소치에서도 넘치면 경고."""
    if not lines:
        return base_pt
    pt = base_pt
    while pt > min_pt:
        if _needed_h(lines, box_w, pt, line, space_after) <= box_h:
            return pt
        pt -= 0.5
    if _needed_h(lines, box_w, min_pt, line, space_after) > box_h and label:
        _log(f"[deck] 넘침 경고: {label} — {min_pt}pt 로도 박스를 초과")
    return min_pt


def _fit_each(items: List[str], box_w: float, box_h: float, base_pt: float,
              line: float = 1.3, min_pt: float = 10.5, label: str = "") -> float:
    """항목마다 '자기 박스'를 갖는 레이아웃용 — 가장 빡빡한 항목에 맞춘 공통 크기.

    _fit() 은 여러 문단이 한 박스에 쌓이는 경우(불릿 목록)를 가정하므로,
    목차·보기·표 셀처럼 항목별 박스가 따로인 곳에 쓰면 높이를 과대 계산한다.
    """
    vals = [str(t) for t in items if str(t).strip()]
    if not vals:
        return base_pt
    return min(_fit([v], box_w, box_h, base_pt, line, 0.0, min_pt, label)
               for v in vals)


def _trim_to_fit(lines: List[str], box_w: float, box_h: float, pt: float,
                 line: float, space_after: float, label: str = "") -> List[str]:
    """최소 크기에서도 안 들어가면 뒤에서부터 항목을 덜어낸다(조용히 자르지 않고 경고)."""
    out = list(lines)
    while len(out) > 1 and _needed_h(out, box_w, pt, line, space_after) > box_h:
        dropped = out.pop()
        _log(f"[deck] 항목 제외: {label} — '{str(dropped)[:24]}…'")
    return out


def _normautofit(tf, scale: float = 1.0, lnspc_red: float = 0.0) -> None:
    """<a:normAutofit> 주입 — PowerPoint 가 편집 시 다시 맞추도록 하는 보조 장치."""
    from pptx.oxml.ns import qn
    try:
        bodyPr = tf._txBody.find(qn("a:bodyPr"))
        if bodyPr is None:
            return
        for tag in ("a:normAutofit", "a:spAutoFit", "a:noAutofit"):
            el = bodyPr.find(qn(tag))
            if el is not None:
                bodyPr.remove(el)
        naf = bodyPr.makeelement(qn("a:normAutofit"), {})
        if scale < 1.0:
            naf.set("fontScale", str(int(scale * 100000)))
        if lnspc_red > 0:
            naf.set("lnSpcReduction", str(int(lnspc_red * 100000)))
        bodyPr.append(naf)
    except Exception:  # noqa: BLE001
        pass


# ── 구성 요소 ─────────────────────────────────────────────────────────────
def add_kicker(slide, text, x=MARGIN, y=KICKER_Y, w=8.0, color=None):
    if not text:
        return
    tf = _text(slide, x, y, w, 0.24, anchor="m")
    _T(_para(tf, True), str(text).upper(), "kicker", color=color)


def add_title(slide, text, *, has_logo=False, y=TITLE_Y, size=None):
    w = 10.6 if has_logo else CONTENT_W
    t = TYPE["title"]
    pt = size or _fit([text or ""], w, TITLE_H, t["size"], t["line"],
                      min_pt=19, label=f"title:{str(text)[:18]}")
    tf = _text(slide, MARGIN, y, w, TITLE_H, anchor="m")
    _T(_para(tf, True, line=t["line"]), text or "", "title", size=pt)


def add_rule(slide, y=RULE_Y, x=MARGIN, w=RULE_W, color="amber"):
    _rrect(slide, x, y, w, RULE_H, fill=color, pill=True)


def add_chip(slide, text, x=MARGIN, y=CHIP_Y, w=CONTENT_W, emphasis=False):
    """한 줄 요약. **글꼴로만** 쓴다 — 반환 = 하단 y.

    예전에는 라운드 배경 + 앰버 점 + 한 문장이었다. 뺀 이유:
      · 알약이 슬라이드마다 반복되니 제목보다 먼저 눈에 들어와 위계가 뒤집혔다
      · 배경색이 사진·도해와 부딪혀 화면이 시끄러웠다
    강조는 색을 채우지 않고 **더 진한 청색**으로 한다.

    ★ 높이(CHIP_H)는 그대로 둔다. 이 함수가 돌려주는 y 를 13종 레이아웃이
      본문 시작점으로 쓰기 때문에, 줄이면 전 레이아웃의 배치가 한꺼번에 틀어진다.
      배경이 없어진 만큼 위아래 여백이 되어 소제목처럼 읽힌다.
    """
    if not text:
        return y
    t = TYPE["chip"]
    # 강조는 색이 아니라 **크기**로 준다. 기본 배색은 navy 가 이미 13.8:1 이라
    # 더 진한 단계가 없어서, 색만으로는 두 템플릿 중 한쪽에서 차이가 안 난다.
    base = t["size"] + (1.5 if emphasis else 0.0)
    pt = _fit([text], w, CHIP_H - 0.16, base, t["line"],
              min_pt=11.5, label=f"chip:{str(text)[:18]}")
    tf = _text(slide, x, y, w, CHIP_H, anchor="m")
    _T(_para(tf, True, line=t["line"]), text, "chip", size=pt,
       color=on_white("navy"))
    return y + CHIP_H


def add_bullets(slide, bullets, x, y, w, h, *, role="bullet", marker=True,
                label=""):
    if not bullets:
        return
    t = TYPE[role]
    items = [str(b) for b in bullets if str(b).strip()]
    if not items:
        return
    sa = 7.0
    # 마커 폭(▸ + 공백)만큼 텍스트 폭이 줄어든다
    eff_w = w - (0.26 if marker else 0.0)
    pt = _fit(items, eff_w, h, t["size"], t["line"], sa, min_pt=11, label=label)
    items = _trim_to_fit(items, eff_w, h, pt, t["line"], sa, label=label)
    tf = _text(slide, x, y, w, h, anchor="t")
    for i, b in enumerate(items):
        p = _para(tf, i == 0, space_after=sa, line=t["line"])
        if marker:
            _hang(p)
            _run(p, "▸ ", pt, "amber", "sb")
        _T(p, b, role, size=pt)
    _normautofit(tf, scale=min(1.0, pt / t["size"]))


def add_photo_slot(slide, x, y, w, h):
    """사진이 아직 없을 때의 **빈 액자**. 자리를 확정해 배치가 두 번 바뀌지 않게 한다.

    글자를 넣지 않는다 — 사진 수집을 건너뛴 덱이 그대로 쓰일 수도 있고, 그때
    '사진 자리' 같은 문구가 남으면 미완성으로 보인다. 대신 옅은 워시와 헤어라인,
    가운데 작은 표식으로 '여기 그림이 들어갈 자리' 를 알린다.
    """
    box = _rrect(slide, x, y, w, h, fill="chip", radius=R_PHOTO,
                 line="navy_12", line_w=1.25)
    # 가운데 작은 원 — 도형이므로 어느 배색에서도 튀지 않는다
    d = min(w, h) * 0.11
    _oval(slide, x + (w - d) / 2, y + (h - d) / 2, d, "navy_30")
    return box


def add_photo(slide, data, x, y, w, h, *, credit=None):
    from pptx.util import Inches
    from pptx.oxml.ns import qn
    try:
        pic = slide.shapes.add_picture(io.BytesIO(data), Inches(x), Inches(y),
                                       Inches(w), Inches(h))
    except Exception:  # noqa: BLE001
        return None
    # 종횡비 맞춰 센터 크롭(가능하면)
    try:
        from PIL import Image
        iw, ih = Image.open(io.BytesIO(data)).size
        tgt, src = w / h, iw / ih
        if src > tgt:
            c = (1 - tgt / src) / 2
            pic.crop_left = c
            pic.crop_right = c
        elif src < tgt:
            c = (1 - src / tgt) / 2
            pic.crop_top = c
            pic.crop_bottom = c
    except Exception:  # noqa: BLE001
        pass
    geom = pic._element.spPr.find(qn("a:prstGeom"))
    if geom is not None:
        geom.set("prst", "roundRect")
        _set_adj(pic, _adj(R_PHOTO, w, h))
    _shadow(pic, blur=0.16, dist=0.04, alpha=14)
    if credit:
        tf = _text(slide, x, y + h + 0.04, w, 0.22, anchor="t")
        _T(_para(tf, True, align="r"), str(credit)[:90], "caption")
    return pic


def add_footer_key(slide, text, y=None):
    """하단 '핵심' 한 줄. 러닝 푸터 위쪽에 놓인다."""
    if not text:
        return
    h = 0.60
    y = BODY_BOTTOM - h if y is None else y
    _rrect(slide, MARGIN, y, CONTENT_W, h, fill="amber_wash", radius=R_CARD)
    _rrect(slide, MARGIN, y, 0.08, h, fill="amber", radius=0.04)
    tf = _text(slide, MARGIN + 0.34, y, CONTENT_W - 0.6, h, anchor="m")
    p = _para(tf, True)
    _run(p, "핵심   ", 12.5, "amber_dk", "sb")
    _run(p, text, 13.5, "navy", "sb")


LOGO_Y, LOGO_H = 0.28, 0.42       # 상단 로고 띠 — 키커(0.52)와 겹치는 높이다


def add_logos(slide, logo_path, logo2_path=None) -> float:
    """상단 로고. 로고1은 **우상단**, 로고2는 **좌상단**.

    반환 = 좌상단 로고가 차지한 폭(in). 0 이면 없다.
    ★ 이 값이 필요한 이유: 키커(0.52~0.76)가 로고 띠(0.28~0.70)와 겹친다.
      좌상단에 로고가 오면 키커를 그만큼 오른쪽으로 밀어야 글자가 안 겹친다.
      로고 폭은 그림 비율에 따라 다르므로 **넣어 보고 재는** 수밖에 없다.

    로고가 없으면 그 자리는 그냥 빈다 — 자리표시를 그리지 않는다.
    """
    from pptx.util import Inches
    left_w = 0.0
    if logo2_path and os.path.isfile(logo2_path):
        try:
            pic = slide.shapes.add_picture(logo2_path, Inches(MARGIN), Inches(LOGO_Y),
                                           height=Inches(LOGO_H))
            left_w = (pic.width or 0) / 914400
        except Exception:  # noqa: BLE001
            pass
    if logo_path and os.path.isfile(logo_path):
        try:
            pic = slide.shapes.add_picture(logo_path, Inches(0), Inches(LOGO_Y),
                                           height=Inches(LOGO_H))
            # 오른쪽 끝을 맞춘다 — 폭이 로고마다 달라 좌표를 박아 두면 삐뚤어진다.
            pic.left = Inches(CONTENT_R - (pic.width or 0) / 914400)
        except Exception:  # noqa: BLE001
            pass
    return left_w


def add_running(slide, footer_text: str, page: Optional[int]):
    """러닝 요소 — 헤어라인 + 좌측 과목·주차 + 우측 페이지 번호."""
    if not footer_text and page is None:
        return
    _rect(slide, MARGIN, FOOT_LINE_Y, CONTENT_W, 0.01, "navy_12")
    if footer_text:
        tf = _text(slide, MARGIN, FOOT_TEXT_Y, 9.0, 0.28, anchor="m")
        _T(_para(tf, True), footer_text, "footer")
    if page is not None:
        tf = _text(slide, CONTENT_R - 1.2, FOOT_TEXT_Y, 1.2, 0.28, anchor="m")
        _T(_para(tf, True, align="r"), str(page), "pageno")


def _head(slide, s, *, has_logo=False, rule=True):
    """상단 존(키커 → 제목 → 룰) 공통 처리. 반환 = 요약줄 시작 y."""
    # 좌상단 로고가 있으면 키커를 그 오른쪽으로 민다(같은 높이라 겹친다).
    add_kicker(slide, s.get("kicker"), x=MARGIN + float(s.get("_kicker_dx") or 0.0))
    add_title(slide, s.get("title", ""), has_logo=has_logo)
    if rule:
        add_rule(slide)
    return CHIP_Y


# ── 사진 배치 프리셋 ───────────────────────────────────────────────────────
# img=(L,T,W,H), tx/tw=텍스트 열, below=이미지가 하단(본문은 위 전폭).
# 사진 슬라이드 '순번'으로 순환한다(절대 인덱스 아님) — 배치 반복 방지.
#
# 세로 기준(첨삭 2026-08-03): 옆단 사진은 **요약줄 윗선(1.88")에 top 을 맞추고**
# 콘텐츠 하단(6.57")까지 내린다. 본문 시작선(2.95")에서 시작하면 칩 옆이 비어
# 위쪽이 헐거워 보인다. 폭만 바꿔 변화를 준다(높이를 줄이면 정렬이 깨진다).
PHOTO_TOP = CHIP_Y          # 1.88 — 요약줄 윗선
PHOTO_BOT = 6.57            # 콘텐츠 하단(러닝 푸터 위)
_PH = PHOTO_BOT - PHOTO_TOP  # 4.69

PHOTO_PRESETS = [
    {"img": (8.55, PHOTO_TOP, 4.18, _PH), "tx": 0.60, "tw": 7.60},  # 우측 표준
    {"img": (0.60, PHOTO_TOP, 4.18, _PH), "tx": 5.13, "tw": 7.60},  # 좌측 표준
    {"img": (7.95, PHOTO_TOP, 4.78, _PH), "tx": 0.60, "tw": 7.00},  # 우측 넓게
    {"img": (0.60, PHOTO_TOP, 4.78, _PH), "tx": 5.72, "tw": 7.01},  # 좌측 넓게
    {"img": (8.20, PHOTO_TOP, 4.53, _PH), "tx": 0.60, "tw": 7.25},  # 우측 중간
]
# '하단 와이드'(2.60, 4.40, 8.13x2.22) 는 뺐다. 비율이 3.7:1 이라 1:1 로 주문한
# 그림을 넣으면 위아래가 잘리고, 본문이 위로 밀려 화면이 비어 보인다.
# below 처리 코드는 남겨 둔다 — 나중에 와이드 자리를 다시 넣을 수 있다.


# ── 타입별 렌더 ───────────────────────────────────────────────────────────
def render_cover(slide, s, img=None):
    """좌측 네이비 밴드 + 앰버 룰 + 우측 타이틀 존."""
    band_w = 4.75
    _rect(slide, 0, 0, band_w, SLIDE_H, "navy")
    _rect(slide, band_w, 0, 0.085, SLIDE_H, "amber")
    _rrect(slide, 0.62, 0.78, 0.60, 0.60, fill="amber", radius=0.12)

    band_note = s.get("band") or s.get("course") or ""
    if band_note:
        tf = _text(slide, 0.62, 5.72, band_w - 1.1, 0.9, anchor="t")
        _T(_para(tf, True), "LECTURE", "kicker_on")
        _T(_para(tf, False, space_after=0, line=1.3), band_note, "cover_sub",
           color="white")

    x = band_w + 0.72
    w = CONTENT_R - x
    if s.get("kicker"):
        add_kicker(slide, s["kicker"], x=x, y=2.42, w=w)
    t = TYPE["cover_title"]
    title = s.get("title", "강의 슬라이드")
    pt = _fit([title], w, 1.95, t["size"], t["line"], min_pt=25,
              label="cover title")
    tf = _text(slide, x, 2.74, w, 1.95, anchor="t")
    _T(_para(tf, True, line=t["line"]), title, "cover_title", size=pt)

    sub = s.get("chip") or s.get("subtitle") or ""
    if sub:
        _rrect(slide, x, 4.86, 0.5, 0.045, fill="amber", pill=True)
        tf2 = _text(slide, x, 5.08, w, 0.9, anchor="t")
        _T(_para(tf2, True, line=TYPE["cover_sub"]["line"]), sub, "cover_sub")


def render_section(slide, s, img=None):
    """구간 전환 — 대형 연번(연한 네이비) + 앰버 바 + 제목."""
    num = str(s.get("num") or "").strip()
    if num:
        tf0 = _text(slide, CONTENT_R - 4.2, 1.55, 4.2, 2.0, anchor="m")
        _T(_para(tf0, True, align="r"), num.zfill(2), "sec_num")
    _rrect(slide, MARGIN, 3.12, 0.11, 1.26, fill="amber", pill=True)
    t = TYPE["sec_title"]
    title = s.get("title", "")
    pt = _fit([title], CONTENT_W - 0.6, 1.5, t["size"], t["line"], min_pt=22,
              label="section title")
    tf = _text(slide, MARGIN + 0.42, 2.98, CONTENT_W - 0.42, 1.55, anchor="m")
    _T(_para(tf, True, line=t["line"]), title, "sec_title", size=pt)
    if s.get("chip"):
        tf2 = _text(slide, MARGIN + 0.42, 4.60, CONTENT_W - 0.42, 0.9, anchor="t")
        _T(_para(tf2, True, line=1.35), s["chip"], "cover_sub", size=15)


def render_bullets(slide, s, img=None):
    y = _head(slide, s)
    by = add_chip(slide, s.get("chip"), y=y, emphasis=bool(s.get("emphasis")))
    top = max(by + 0.28, BODY_TOP)
    bottom = BODY_BOTTOM
    if s.get("key") or s.get("footer"):
        add_footer_key(slide, s.get("key") or s.get("footer"))
        bottom = BODY_BOTTOM - 0.60 - 0.22
    add_bullets(slide, s.get("bullets", []), MARGIN + 0.04, top,
                CONTENT_W - 0.08, bottom - top,
                label=f"bullets:{str(s.get('title'))[:18]}")


def render_photo(slide, s, img=None):
    """사진 슬라이드. 사진이 없어도 **자리는 남긴다.**

    예전에는 사진이 없으면 전폭 불릿으로 되돌아갔다. 그래서 초안(2단계) 덱에서는
    사진 자리가 보이지 않고, 사진을 넣는 3단계에서야 비로소 자리가 생기며 글이
    한쪽으로 밀렸다 — 배치가 두 번 바뀌는 셈이었다.

    지금은 빈 액자를 그려 자리를 확정한다. 그래서
      · 초안만 보고도 최종 배치를 알 수 있다
      · 사진 수집(3단계)을 건너뛰고 대본·영상으로 바로 가도 배치가 그대로다
      · 직접 그림을 그려 넣을 때 어느 크기·위치에 맞출지 눈으로 보인다
    """
    y = _head(slide, s)
    emph = bool(s.get("emphasis"))
    p = PHOTO_PRESETS[s.get("_photo_ord", 0) % len(PHOTO_PRESETS)]
    if img:
        add_photo(slide, img, *p["img"], credit=s.get("_credit_short"))
    else:
        add_photo_slot(slide, *p["img"])
    tx, tw = p["tx"], p["tw"]
    by = add_chip(slide, s.get("chip"), x=tx, y=y, w=tw, emphasis=emph)
    # 칩이 없으면 본문을 사진 윗선에 맞춘다(칩이 있을 때는 칩 아래).
    top = max(by + 0.26, BODY_TOP) if s.get("chip") else PHOTO_TOP
    bottom = (p["img"][1] - 0.22) if p.get("below") else BODY_BOTTOM
    add_bullets(slide, s.get("bullets", []), tx + 0.04, top, tw - 0.08,
                max(0.8, bottom - top),
                label=f"photo:{str(s.get('title'))[:18]}")


def render_process(slide, s, img=None):
    y = _head(slide, s)
    add_chip(slide, s.get("chip"), y=y)
    items = s.get("items") or [{"label": b} for b in s.get("bullets", [])]
    items = items[:5] or [{"label": "항목"}]
    n = len(items)
    gap = 0.36
    node_h = 1.42
    node_w = (CONTENT_W - gap * (n - 1)) / n
    ny = BODY_TOP + 0.18
    has_desc = any(it.get("desc") for it in items)
    lt, ld = TYPE["node_label"], TYPE["node_desc"]

    labels = [str(it.get("label", "")) for it in items]
    lpt = min(_fit([l], node_w - 0.24, node_h - 0.2, lt["size"], lt["line"],
                   min_pt=10.5, label="process label") for l in labels)
    descs = [str(it.get("desc") or "") for it in items]
    dpt = ld["size"]
    if has_desc:
        dpt = min(_fit([d], node_w - 0.14, 1.45, ld["size"], ld["line"],
                       min_pt=9, label="process desc") for d in descs if d)

    for i, it in enumerate(items):
        x = MARGIN + i * (node_w + gap)
        # 단일 계조: 마지막 노드만 앰버로 도착점 표시
        last = (i == n - 1) and n > 1
        fill = "amber" if last else _RAMP[min(i, len(_RAMP) - 1)]
        txt_color = "navy" if last else "white"
        node = _rrect(slide, x, ny, node_w, node_h, fill=fill, radius=R_CARD)
        _shadow(node, blur=0.12, dist=0.03, alpha=10)
        tf = _text(slide, x + 0.12, ny, node_w - 0.24, node_h, anchor="m")
        _T(_para(tf, True, align="c", line=lt["line"]), it.get("label", ""),
           "node_label", size=lpt, color=txt_color)
        if it.get("desc"):
            tf2 = _text(slide, x + 0.07, ny + node_h + 0.14, node_w - 0.14, 1.45)
            _T(_para(tf2, True, align="c", line=ld["line"]), it["desc"],
               "node_desc", size=dpt)
        if i < n - 1:
            _arrow(slide, x + node_w + 0.05, ny + node_h / 2 - 0.13,
                   gap - 0.10, 0.26, "amber")
    add_footer_key(slide, s.get("key") or s.get("footer"))


def render_cards(slide, s, img=None):
    y = _head(slide, s)
    add_chip(slide, s.get("chip"), y=y)
    items = s.get("items") or [{"label": b} for b in s.get("bullets", [])]
    items = items[:4] or [{"label": "항목"}]
    numbered = bool(s.get("numbered", False))   # 비순차 분류는 숫자 없음
    n = len(items)
    gap = 0.34
    card_w = (CONTENT_W - gap * (n - 1)) / n
    cy = BODY_TOP + 0.05
    has_key = bool(s.get("key") or s.get("footer"))
    card_h = (BODY_BOTTOM - 0.82 if has_key else BODY_BOTTOM) - cy

    lt, dt = TYPE["card_label"], TYPE["card_desc"]
    labels = [str(it.get("label", "")) for it in items]
    lpt = min(_fit([l], card_w - 0.62, 0.62, lt["size"], lt["line"],
                   min_pt=11.5, label="card label") for l in labels)
    descs = [str(it.get("desc") or "") for it in items]
    label_bottom = 1.16 if numbered else 0.98
    desc_h = card_h - label_bottom - 0.26
    dpt = dt["size"]
    if any(descs):
        dpt = min(_fit([d], card_w - 0.62, desc_h, dt["size"], dt["line"],
                       min_pt=9.5, label="card desc") for d in descs if d)

    for i, it in enumerate(items):
        x = MARGIN + i * (card_w + gap)
        accent = _RAMP[min(i, len(_RAMP) - 1)]
        card = _rrect(slide, x, cy, card_w, card_h, fill="white", radius=R_CARD,
                      line="navy_12", line_w=1.0)
        _shadow(card, blur=0.13, dist=0.03, alpha=9)
        if numbered:
            _oval(slide, x + 0.30, cy + 0.30, 0.52, accent)
            tfb = _text(slide, x + 0.30, cy + 0.30, 0.52, 0.52, anchor="m")
            _T(_para(tfb, True, align="c"), str(i + 1), "card_label", size=15,
               color="white")
        else:
            _rrect(slide, x, cy, card_w, 0.14, fill=accent, radius=0.07)
            _rrect(slide, x + 0.30, cy + 0.42, 0.30, 0.075, fill=accent, pill=True)
        tft = _text(slide, x + 0.30, cy + label_bottom - 0.62, card_w - 0.60, 0.62,
                    anchor="m")
        _T(_para(tft, True, line=lt["line"]), it.get("label", ""), "card_label",
           size=lpt, color="navy")
        if it.get("desc"):
            tfd = _text(slide, x + 0.30, cy + label_bottom, card_w - 0.60, desc_h)
            _T(_para(tfd, True, line=dt["line"]), it["desc"], "card_desc", size=dpt)
    add_footer_key(slide, s.get("key") or s.get("footer"))


def render_compare(slide, s, img=None):
    y = _head(slide, s)
    add_chip(slide, s.get("chip"), y=y)
    items = (s.get("items") or [])[:2]
    while len(items) < 2:
        items.append({"label": "", "lines": []})
    cy = BODY_TOP + 0.05
    col_h = BODY_BOTTOM - cy
    col_w = (CONTENT_W - 0.42) / 2
    head_h = 0.60
    # 밴드에 글자가 얹히므로 대비가 모자라면 진한 짝으로 바꾼다(템플릿마다 다르다)
    fills = (band("navy"), band("amber"))
    # ★ 글자색을 고정하지 않고 **배경 밝기로 고른다.** 색을 박아 두면 템플릿에서
    #   밴드 색을 바꿨을 때 대비가 무너진다(초록 밴드에 남색 글씨가 흐릿했다).
    heads = tuple(on_color(f) for f in fills)
    td = TYPE["td"]

    for i, it in enumerate(items):
        x = MARGIN + i * (col_w + 0.42)
        col = _rrect(slide, x, cy, col_w, col_h, fill="chip", radius=R_CARD)
        _shadow(col, blur=0.12, dist=0.03, alpha=8)
        _rrect(slide, x, cy, col_w, head_h, fill=fills[i], radius=R_CARD)
        _rect(slide, x, cy + head_h - 0.10, col_w, 0.10, fills[i])
        tft = _text(slide, x + 0.2, cy, col_w - 0.4, head_h, anchor="m")
        _T(_para(tft, True, align="c"), it.get("label", ""), "card_label",
           size=15, color=heads[i])
        lines = it.get("lines") or ([it["desc"]] if it.get("desc") else [])
        lines = [str(l) for l in lines if str(l).strip()]
        if not lines:
            continue
        # 라벨만 가운데, 본문은 좌측 정렬(문장이 가운데면 읽히지 않는다)
        box_w, box_h = col_w - 0.66, col_h - head_h - 0.42
        sa = 6.0
        pt = _fit(lines, box_w, box_h, 14.0, td["line"], sa, min_pt=10,
                  label=f"compare:{it.get('label')}")
        lines = _trim_to_fit(lines, box_w, box_h, pt, td["line"], sa,
                             label=f"compare:{it.get('label')}")
        tf = _text(slide, x + 0.34, cy + head_h + 0.26, box_w, box_h, anchor="t")
        for j, ln in enumerate(lines):
            p = _para(tf, j == 0, space_after=sa, line=td["line"])
            _hang(p, 0.22)
            _run(p, "· ", pt, "amber", "sb")
            _T(p, ln, "td", size=pt)


def _cell_border(cell, color="white", w_pt=1.0):
    from pptx.oxml.ns import qn
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for tag in ("a:lnL", "a:lnR", "a:lnT", "a:lnB"):
        for el in tcPr.findall(qn(tag)):
            tcPr.remove(el)
    # lnL, lnR, lnT, lnB 는 tcPr 의 맨 앞 순서
    for tag in ("a:lnB", "a:lnT", "a:lnR", "a:lnL"):
        ln = tcPr.makeelement(qn(tag), {"w": str(int(w_pt * 12700)),
                                        "cap": "flat", "cmpd": "sng",
                                        "algn": "ctr"})
        fill = ln.makeelement(qn("a:solidFill"), {})
        fill.append(fill.makeelement(qn("a:srgbClr"), {"val": _hexstr(color)}))
        ln.append(fill)
        tcPr.insert(0, ln)


def render_table(slide, s, img=None):
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.oxml.ns import qn

    y0 = _head(slide, s)
    by = add_chip(slide, s.get("chip"), y=y0)
    rows_data = s.get("rows")
    if not rows_data:
        return render_bullets(slide, s)
    headers = s.get("headers") or []
    y = max(by + 0.30, BODY_TOP)
    ncol = len(headers) or max(len(r) for r in rows_data)
    nrow = len(rows_data) + (1 if headers else 0)
    avail = BODY_BOTTOM - y
    row_h = min(0.46, max(0.30, avail / max(nrow, 1)))
    tbl_shape = slide.shapes.add_table(nrow, ncol, Inches(MARGIN), Inches(y),
                                       Inches(CONTENT_W), Inches(row_h * nrow))
    tbl = tbl_shape.table

    # PowerPoint 기본 파랑 밴딩 스타일 제거 — 셀 채움을 직접 지정한다
    try:
        tbl.first_row = bool(headers)
        tbl.horz_banding = False
        tbl.vert_banding = False
        tblPr = tbl._tbl.find(qn("a:tblPr"))
        if tblPr is not None:
            for el in tblPr.findall(qn("a:tableStyleId")):
                tblPr.remove(el)
    except Exception:  # noqa: BLE001
        pass

    # 문장 열은 좌측, 짧은 값 열은 가운데
    col_len = []
    for c in range(ncol):
        vals = [str(r[c]) for r in rows_data if c < len(r)]
        col_len.append(max((len(v) for v in vals), default=0))
    aligns = [PP_ALIGN.LEFT if L > 12 else PP_ALIGN.CENTER for L in col_len]

    th, td = TYPE["th"], TYPE["td"]
    body_pt = td["size"]
    if rows_data:
        cell_w = CONTENT_W / ncol
        body_pt = min(
            _fit_each([str(r[c]) for r in rows_data if c < len(r)],
                      cell_w - 0.24, row_h - 0.10, td["size"], 1.16,
                      min_pt=8.5, label=f"table col{c}")
            for c in range(ncol))

    def _style(cell, text, *, header, align, bg):
        cell.text = str(text)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = cell.margin_right = Inches(0.10)
        cell.margin_top = cell.margin_bottom = Inches(0.03)
        cell.fill.solid()
        cell.fill.fore_color.rgb = _rgb(bg)
        _cell_border(cell, "white", 1.0)
        for para in cell.text_frame.paragraphs:
            para.alignment = align
            para.line_spacing = 1.16
            for run in para.runs:
                run.font.size = Pt(th["size"] if header else body_pt)
                run.font.bold = _FS.is_bold("sb" if header else "r")
                run.font.color.rgb = _rgb("white" if header else "ink")
                pptx_font.apply_face(run, _FS.face("sb" if header else "r"))

    r0 = 0
    if headers:
        for c, htxt in enumerate(headers):
            _style(tbl.cell(0, c), htxt, header=True, align=PP_ALIGN.CENTER,
                   bg="navy")
        r0 = 1
    for ri, row in enumerate(rows_data):
        bg = "white" if ri % 2 == 0 else "chip"
        for c in range(ncol):
            _style(tbl.cell(ri + r0, c), row[c] if c < len(row) else "",
                   header=False, align=aligns[c], bg=bg)


def render_quiz(slide, s, img=None):
    """형성평가 1문항 — 문제(네이비 칩) + 보기 알약 + 정답 뱃지 우하단."""
    s = dict(s)
    s.setdefault("kicker", "형성평가")
    y = _head(slide, s)
    question = s.get("question") or s.get("chip") or ""
    by = add_chip(slide, question, y=y, emphasis=True) if question else y

    choices = s.get("choices") or s.get("items") or s.get("bullets") or []
    choices = [str(c).strip() for c in choices if str(c).strip()][:5]
    top = max(by + 0.34, BODY_TOP)
    answer = s.get("answer") or s.get("key")
    bottom = BODY_BOTTOM - (0.72 if answer else 0.0)

    if choices:
        marks = "①②③④⑤"
        n = len(choices)
        gap = 0.16
        row_h = min(0.82, max(0.52, (bottom - top - gap * (n - 1)) / n))
        pt = _fit_each(choices, CONTENT_W - 1.2, row_h - 0.12, 15.0, 1.2,
                       min_pt=11, label="quiz choice")
        for i, c in enumerate(choices):
            ry = top + i * (row_h + gap)
            _rrect(slide, MARGIN, ry, CONTENT_W, row_h, fill="chip", pill=True)
            _oval(slide, MARGIN + 0.20, ry + row_h / 2 - 0.15, 0.30, "navy")
            tfm = _text(slide, MARGIN + 0.20, ry, 0.30, row_h, anchor="m")
            _T(_para(tfm, True, align="c"), marks[i], "card_label", size=12,
               color="white")
            tf = _text(slide, MARGIN + 0.66, ry, CONTENT_W - 0.9, row_h, anchor="m")
            _T(_para(tf, True, line=1.2), c, "bullet", size=pt)
    else:
        add_bullets(slide, s.get("bullets", []), MARGIN + 0.04, top,
                    CONTENT_W - 0.08, bottom - top, label="quiz body")

    if answer:
        w = min(3.6, max(1.5, 0.30 + _measure(f"정답  {answer}", 13.0)))
        ay = BODY_BOTTOM - 0.50
        _rrect(slide, CONTENT_R - w, ay, w, 0.50, fill="amber", pill=True)
        tf = _text(slide, CONTENT_R - w + 0.16, ay, w - 0.32, 0.50, anchor="m")
        p = _para(tf, True, align="c")
        _run(p, "정답  ", 12.0, "amber_dk", "sb")
        _T(p, answer, "answer")


def render_objectives(slide, s, img=None):
    """학습목표 — 대형 연번 + 목표 문장. 이미지 없음."""
    s = dict(s)
    s.setdefault("kicker", "학습목표")
    y = _head(slide, s)
    by = add_chip(slide, s.get("chip"), y=y) if s.get("chip") else y
    goals = s.get("items") or s.get("bullets") or []
    goals = [(g.get("label") if isinstance(g, dict) else str(g)) for g in goals]
    goals = [str(g).strip() for g in goals if str(g).strip()][:4]
    if not goals:
        return render_bullets(slide, s)
    top = max(by + 0.34, BODY_TOP)
    n = len(goals)
    gap = 0.20
    row_h = min(1.30, (BODY_BOTTOM - top - gap * (n - 1)) / n)
    pt = _fit_each(goals, CONTENT_W - 1.72, row_h - 0.14, 17.0, 1.34,
                   min_pt=12, label="objective")
    for i, g in enumerate(goals):
        ry = top + i * (row_h + gap)
        tfn = _text(slide, MARGIN, ry, 1.30, row_h, anchor="m")
        _T(_para(tfn, True), str(i + 1).zfill(2), "big_num")
        _rect(slide, MARGIN + 1.16, ry + 0.10, 0.02, row_h - 0.20, "navy_12")
        tf = _text(slide, MARGIN + 1.42, ry, CONTENT_W - 1.42, row_h, anchor="m")
        _T(_para(tf, True, line=1.34), g, "bullet", size=pt)


def render_agenda(slide, s, img=None):
    """목차 — 2열 번호 목록. active 로 현재 구간 강조."""
    s = dict(s)
    s.setdefault("kicker", "목차")
    y = _head(slide, s)
    by = add_chip(slide, s.get("chip"), y=y) if s.get("chip") else y
    items = s.get("items") or s.get("bullets") or []
    items = [(i.get("label") if isinstance(i, dict) else str(i)) for i in items]
    items = [str(i).strip() for i in items if str(i).strip()][:10]
    if not items:
        return render_bullets(slide, s)
    active = s.get("active")
    top = max(by + 0.34, BODY_TOP)
    ncol = 2 if len(items) > 4 else 1
    per = math.ceil(len(items) / ncol)
    col_w = (CONTENT_W - 0.5) / ncol
    row_h = min(0.80, (BODY_BOTTOM - top) / per)
    pt = _fit_each(items, col_w - 0.85, row_h - 0.12, 15.5, 1.25, min_pt=11,
                   label="agenda")
    for i, it in enumerate(items):
        c, r = divmod(i, per)
        x = MARGIN + c * (col_w + 0.5)
        ry = top + r * row_h
        on = (active is not None and int(active) == i + 1)
        if on:
            _rrect(slide, x - 0.10, ry + 0.04, col_w + 0.10, row_h - 0.08,
                   fill="chip", radius=R_CARD)
        _oval(slide, x, ry + row_h / 2 - 0.15, 0.30, "navy" if on else "navy_30")
        tfn = _text(slide, x, ry, 0.30, row_h, anchor="m")
        _T(_para(tfn, True, align="c"), str(i + 1), "card_label", size=11.5,
           color="white")
        tf = _text(slide, x + 0.46, ry, col_w - 0.5, row_h, anchor="m")
        _T(_para(tf, True, line=1.25), it, "bullet", size=pt,
           color=("navy" if on else "ink"), weight=("sb" if on else "r"))


def render_closing(slide, s, img=None):
    """마무리 — 요약 / 다음 차시 / 과제 3블록."""
    s = dict(s)
    s.setdefault("kicker", "마무리")
    y = _head(slide, s)
    add_chip(slide, s.get("chip"), y=y)
    items = s.get("items") or []
    if not items:
        bl = [str(b) for b in s.get("bullets", []) if str(b).strip()]
        if not bl:
            return render_bullets(slide, s)
        labels = ["오늘의 요약", "다음 차시", "과제"]
        per = math.ceil(len(bl) / 3) or 1
        items = [{"label": labels[i], "lines": bl[i * per:(i + 1) * per]}
                 for i in range(3) if bl[i * per:(i + 1) * per]]
    items = items[:3]
    n = len(items)
    gap = 0.34
    col_w = (CONTENT_W - gap * (n - 1)) / n
    cy = BODY_TOP + 0.05
    col_h = BODY_BOTTOM - cy
    td = TYPE["td"]
    for i, it in enumerate(items):
        x = MARGIN + i * (col_w + gap)
        box = _rrect(slide, x, cy, col_w, col_h, fill="white", radius=R_CARD,
                     line="navy_12", line_w=1.0)
        _shadow(box, blur=0.13, dist=0.03, alpha=9)
        _rrect(slide, x, cy, col_w, 0.14, fill=_RAMP[min(i, 3)], radius=0.07)
        tfl = _text(slide, x + 0.28, cy + 0.34, col_w - 0.56, 0.5, anchor="m")
        _T(_para(tfl, True), it.get("label", ""), "card_label")
        lines = it.get("lines") or ([it["desc"]] if it.get("desc") else [])
        lines = [str(l) for l in lines if str(l).strip()]
        if not lines:
            continue
        box_w, box_h = col_w - 0.56, col_h - 1.10
        sa = 6.0
        pt = _fit(lines, box_w, box_h, 13.0, td["line"], sa, min_pt=9.5,
                  label=f"closing:{it.get('label')}")
        lines = _trim_to_fit(lines, box_w, box_h, pt, td["line"], sa,
                             label=f"closing:{it.get('label')}")
        tf = _text(slide, x + 0.28, cy + 0.94, box_w, box_h, anchor="t")
        for j, ln in enumerate(lines):
            p = _para(tf, j == 0, space_after=sa, line=td["line"])
            _hang(p, 0.22)
            _run(p, "· ", pt, "amber", "sb")
            _T(p, ln, "td", size=pt)


def render_stat(slide, s, img=None):
    """수치 강조 — 대형 숫자 + 단위 + 한 줄 설명(최대 3개)."""
    y = _head(slide, s)
    add_chip(slide, s.get("chip"), y=y)
    items = s.get("items") or []
    if not items:
        return render_bullets(slide, s)
    items = items[:3]
    n = len(items)
    gap = 0.34
    col_w = (CONTENT_W - gap * (n - 1)) / n
    cy = BODY_TOP + 0.15
    has_key = bool(s.get("key") or s.get("footer"))
    col_h = (BODY_BOTTOM - 0.82 if has_key else BODY_BOTTOM) - cy
    for i, it in enumerate(items):
        x = MARGIN + i * (col_w + gap)
        _rrect(slide, x, cy, col_w, col_h, fill="chip", radius=R_CARD)
        _rrect(slide, x + col_w / 2 - 0.28, cy + 0.34, 0.56, 0.05, fill="amber",
               pill=True)
        num = str(it.get("label", ""))
        npt = _fit([num], col_w - 0.5, 1.3, TYPE["stat_num"]["size"], 1.1,
                   min_pt=26, label="stat num")
        tfn = _text(slide, x + 0.2, cy + 0.72, col_w - 0.4, 1.4, anchor="m")
        _T(_para(tfn, True, align="c"), num, "stat_num", size=npt)
        if it.get("desc"):
            dh = col_h - 2.3
            dpt = _fit([str(it["desc"])], col_w - 0.7, dh, 13.5, 1.35,
                       min_pt=10, label="stat desc")
            tfd = _text(slide, x + 0.35, cy + 2.15, col_w - 0.7, dh, anchor="t")
            _T(_para(tfd, True, align="c", line=1.35), it["desc"], "td", size=dpt)
    add_footer_key(slide, s.get("key") or s.get("footer"))


_RENDER = {
    "cover": render_cover,
    "section": render_section,
    "photo": render_photo,
    "process": render_process,
    "cards": render_cards,
    "compare": render_compare,
    "table": render_table,
    "bullets": render_bullets,
    "quiz": render_quiz,
    "objectives": render_objectives,
    "agenda": render_agenda,
    "closing": render_closing,
    "stat": render_stat,
}

# 러닝 푸터·페이지 번호를 넣지 않는 타입
_NO_RUNNING = ("cover", "section")


def build_deck(plan: List[Dict], template_path: Optional[str] = None,
               images: Optional[Dict[int, bytes]] = None,
               deck_title: str = "강의 슬라이드",
               logo_path: Optional[str] = None, *,
               logo2_path: Optional[str] = None,
               footer: str = "", assets_dir=None,
               embed_font: bool = True,
               template: Optional[str] = None) -> Optional[bytes]:
    """슬라이드플랜 → 디자인된 .pptx 바이트. python-pptx 미설치 시 None.

    embed_font=False 면 Pretendard 를 '이름으로만' 지정한다(파일 ~3MB 작아지지만
    폰트가 설치되지 않은 PC 에서는 대체 폰트로 열린다).

    template — templates/<이름>.json 의 배색·글꼴을 적용한다. None 이면 기본 배색.
    """
    global _FS
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except Exception:  # noqa: BLE001
        return None

    # ★ 색·글꼴을 먼저 갈아끼운다. 도형을 그리기 시작한 뒤에 바꾸면 앞뒤 장의
    #   색이 섞인다.
    from core import palette as _pal
    _tpl = _pal.apply(template)

    adir = assets_dir or ASSETS_DIR
    _FS = pptx_font.font_set(adir)
    # 템플릿이 **Pretendard 아닌** 글꼴을 지정했을 때만 갈아끼운다.
    #   · Pretendard 는 font_set() 이 이미 다룬다 — assets 에 파일이 있으면 임베드하고,
    #     없으면 맑은 고딕으로 폴백한다. 여기서 덮어쓰면 그 폴백이 깨져서
    #     설치도 안 된 Pretendard 를 이름으로만 지정하게 된다(실측: 회귀했다).
    #   · 임베드는 Pretendard 만 가능하므로 다른 글꼴은 이름만 넣는다.
    _tpl_font = (_tpl.get("font") or "").strip()
    if _tpl_font and _tpl_font.lower() != pptx_font.PREFERRED.lower():
        _FS = pptx_font.FontSet(regular=_tpl_font, semibold=_tpl_font,
                                bold=_tpl_font, embedded=False)
        _log(f"[font] 템플릿 글꼴: {_tpl_font} (임베드 불가 — 이름만 지정)")

    images = images or {}
    use_tpl = bool(template_path and os.path.isfile(template_path))
    if use_tpl:
        try:
            prs = Presentation(template_path)
            _remove_all_slides(prs)
        except Exception:  # noqa: BLE001
            prs = Presentation()
            use_tpl = False
    else:
        prs = Presentation()
    if not use_tpl:
        prs.slide_width = Inches(SLIDE_W)
        prs.slide_height = Inches(SLIDE_H)

    if _FS.embedded and embed_font:
        got = pptx_font.embed_fonts(prs, adir)
        if not got:
            _FS = pptx_font.SYSTEM_SET     # 임베드 실패 → 안전하게 시스템 폰트로
            _log("[font] 임베드 불가 — 맑은 고딕으로 폴백")
        else:
            _log(f"[font] 임베드: {', '.join(got)}")
    elif _FS.embedded:
        _log("[font] 임베드 끄기 — Pretendard 이름만 지정")

    has_logo = bool(logo_path and os.path.isfile(logo_path))
    has_logo2 = bool(logo2_path and os.path.isfile(logo2_path))
    layout = _blank_layout(prs)
    photo_ord = 0
    sec_num = 0
    fallbacks: List[int] = []

    for i, s in enumerate(plan):
        slide = prs.slides.add_slide(layout)
        _strip_placeholders(slide)
        typ = (s.get("type") or "bullets").lower()
        s["_idx"] = i
        if typ == "photo":
            s["_photo_ord"] = photo_ord
            photo_ord += 1
        if typ == "section":
            sec_num += 1
            s.setdefault("num", sec_num)
        s["_has_logo"] = has_logo
        # ★ 로고를 **먼저** 놓는다. 좌상단 로고의 실제 폭을 알아야 키커를 밀 수 있고,
        #   폭은 그림을 넣어 봐야 안다. 표지는 자체 레이아웃이라 넣지 않는다.
        s["_kicker_dx"] = 0.0
        if typ != "cover" and (has_logo or has_logo2):
            lw = add_logos(slide, logo_path if has_logo else None,
                           logo2_path if has_logo2 else None)
            s["_kicker_dx"] = (lw + 0.22) if lw else 0.0
        fn = _RENDER.get(typ, render_bullets)
        try:
            fn(slide, s, images.get(i))
        except Exception as e:  # noqa: BLE001
            _log(f"[deck] slide {i + 1} ({typ}) 렌더 오류: {e}")
            fallbacks.append(i + 1)
            try:
                render_bullets(slide, s)
            except Exception:  # noqa: BLE001
                pass
        if typ not in _NO_RUNNING:
            add_running(slide, footer, i + 1)

    if fallbacks:
        _log(f"[deck] 폴백 렌더된 슬라이드: {fallbacks}")

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ── 아트디렉터(LLM) 패스 — 개요 md → 슬라이드플랜 JSON ─────────────────────
_ART_SYS = (
    "너는 강의 슬라이드 아트디렉터다. 주어진 슬라이드 개요를 슬라이드플랜 JSON 배열로만 변환한다. "
    "설명·머리말·코드펜스 없이 JSON 배열만 출력한다."
)

_ART_RULES = """개요의 '### 슬라이드 N' 블록 하나당 JSON 객체 하나를 순서대로 만든다(개수 유지).
각 객체 스키마:
{
  "type": "section|photo|process|cards|compare|table|quiz|objectives|agenda|closing|stat|bullets",
  "title": "간결·정확한 제목",
  "kicker": "상단 초소형 라벨(선택, 6자 이내. 예 '개념'·'사례'·'이론')",
  "chip": "요약 한 문장(명사형). 없으면 생략 가능",
  "bullets": ["불릿 3~5개(명사형 어미)"],
  "items": [{"label":"짧은 이름","desc":"한 줄 설명"}],
  "image_query": "photo형일 때만, 영어 검색어",
  "key": "하단 핵심 한 줄(선택, 명사형)",
  "level": "기억|이해|적용|분석|평가|창조 — 이 슬라이드가 지지하는 차시 학습목표의 인지수준",
  "numbered": false,
  "emphasis": false
}
인지수준(level) — 반드시 채운다:
- 개요의 각 블록에 있는 **'지지 목표: (인지수준) 목표N'** 줄을 그대로 옮긴다.
  그 줄이 없으면 개요 맨 앞 '차시 학습목표'의 인지수준 태그([기억]…[창조])를 보고
  이 슬라이드가 어느 목표를 지지하는지 판단해 적는다. 새 목표를 만들지 않는다.
- '지지 목표' 줄은 설계용 메타이므로 bullets 에 넣지 않는다.
- level 은 슬라이드에 인쇄되지 않는다(학습자용 화면에 설계 태그를 노출하지 않는다).
  교수자가 목표–자료 정렬을 점검하는 데 쓴다.
**인지수준 → 레이아웃 정렬(구성적 정렬 원칙)**:
- 기억: cards(용어·항목 병렬) · bullets · table(정의 목록)
- 이해: photo(예시·사례로 의미 부여) · bullets · process(흐름 설명)
- 적용: process(절차·단계) · cards(적용 조건) · quiz
- 분석: compare(두 개념 대비) · table(기준별 비교) · stat(수치 근거)
- 평가: compare(장단·타당성) · table(판단 기준) · quiz
- 창조: process(산출 절차) · cards(설계 요소) · closing(과제로 연결)
목표가 분석·평가인데 레이아웃이 단순 나열(bullets)이면 정렬 오류다 — 비교·표로 바꾼다.
문체·제목:
- 불릿·라벨은 **명사형 어미**(예 "즉시 피드백 제공" → "즉시 피드백"). 존댓말·서술체·교수자 나레이션 문장 금지.
- 제목은 간결·일관("정의와 과정"·"개요"·"논의" 같은 군더더기 제거).
- **불릿 한 줄은 32자 이내.** 길면 두 항목으로 쪼갠다(슬라이드가 넘친다).
**사진 비중 — 반드시 지킬 것**:
- 이 덱은 강의 **영상**으로 렌더된다. 도형만 이어지면 화면이 단조로워진다.
- **본문 슬라이드의 약 절반을 "photo" 로 배정한다.** 개념 설명·사례·인물·현장
  장면처럼 시각화가 도움이 되는 것은 photo 를 우선한다.
- photo 로 하면 `image_query`(영어 검색어)를 **반드시** 넣는다. 없으면 사진을 못 찾는다.
- 표·수치·순서·2단 대비처럼 **구조가 정보 자체인 것만** table/stat/process/compare 로 둔다.
- 학습목표·목차·퀴즈·마무리에는 사진을 넣지 않는다(objectives/agenda/quiz/closing).
타입 선택 규칙:
- "cover"는 만들지 마라(표지는 시스템이 맨 앞에 자동 추가).
- **학습목표 슬라이드 → "objectives"** (bullets 에 목표 2~3개). 이미지 없음.
- **목차·구성 안내 → "agenda"** (items[].label 또는 bullets).
- **퀴즈·형성평가 1문항 → "quiz"**: {"question":"문제", "choices":["보기1",…4개], "answer":"③"}.
  한 슬라이드에 한 문제. 정답은 시스템이 우하단 뱃지로 배치한다.
- **마무리·요약 → "closing"**: items 3개 {"label":"오늘의 요약|다음 차시|과제","lines":[…]}.
- 순서·단계·절차가 핵심: "process" (items[].label 2~5개). 화살표로 연결됨.
- 병렬·분류·구성요소: "cards" (items[].label + desc, 2~4개). **비순차이므로 numbered=false(기본).**
- 두 개념 대비/비교: "compare" (items 2개, 각 {"label","lines":[...]}).
- 표가 자연스러운 것: "table" (headers:[...], rows:[[...]]). 5열·8행 이내.
- 수치가 메시지인 것: "stat" (items 2~3개, {"label":"85%","desc":"설명"}).
- 도입·구간 전환: "section".
- 사진이 이해를 돕는 개념/사례/인물: "photo" + 구체적 영어 image_query(막연한 'education' 금지).
  예) 행동주의→'Pavlov conditioning experiment', 인지주의→'human brain memory diagram',
  매체→'classroom projector lesson', 인물→'portrait of a scholar'. 전체의 약 40~55%.
이미지 배치 규칙(중요):
- **objectives·quiz·closing·table·순수 정리 슬라이드는 photo 로 하지 말 것(이미지 없음).**
- 사례·실물·인물·대표 이론 슬라이드에는 photo 지정.
강조(emphasis):
- **emphasis=true 는 문제/질문 칩에만.** quiz 는 자동 강조되므로 따로 켜지 않는다.
공통:
- 같은 타입이 3장 이상 연속되지 않게 섞는다. 동일 내용이 반복되면 한 슬라이드로 합친다.
- process/cards/stat 엔 key(하단 핵심 한 줄, 명사형)를 넣어도 좋다.
"""


def _extract_json_array(text: str):
    """JSON 배열 파싱. 잘린 응답도 완결된 {…} 객체만 골라 살려낸다."""
    if not text:
        return None
    t = re.sub(r"^```[a-zA-Z]*\s*|\s*```$", "", text.strip())
    i = t.find("[")
    if i == -1:
        return None
    j = t.rfind("]")
    if j > i:
        try:
            return json.loads(t[i:j + 1])
        except Exception:  # noqa: BLE001
            pass
    # 살리기: 배열 안의 완결된 최상위 { … } 객체들을 순서대로 파싱
    objs, depth, start, instr, esc = [], 0, None, False, False
    for k in range(i + 1, len(t)):
        c = t[k]
        if instr:
            if esc:
                esc = False
            elif c == "\\":
                esc = True
            elif c == '"':
                instr = False
            continue
        if c == '"':
            instr = True
        elif c == "{":
            if depth == 0:
                start = k
            depth += 1
        elif c == "}":
            if depth > 0:
                depth -= 1
                if depth == 0 and start is not None:
                    try:
                        objs.append(json.loads(t[start:k + 1]))
                    except Exception:  # noqa: BLE001
                        pass
                    start = None
    return objs or None


def _outline_blocks(outline_md: str) -> List[str]:
    """개요 md 를 '### 슬라이드 …' 단위 블록 리스트로 분할."""
    parts = re.split(r"(?m)^(?=#{2,3}\s*슬라이드)", outline_md or "")
    return [p.strip() for p in parts if re.match(r"#{2,3}\s*슬라이드", p.strip())]


# 개요 블록의 '설계 메타' 키 — 슬라이드 본문이 아니므로 불릿으로 넣지 않는다.
# 이걸 걸러내지 않으면 '발표자 노트' 의 나레이션 문단이 통째로 불릿이 된다(실제로 그랬다).
_META_KEYS = ("레이아웃", "시각자료", "강조 신호", "발표자", "노트", "지지",
              "구간", "시간", "슬라이드 수", "이미지", "출처", "비고")
# '- **키**: 값' 형태의 라벨 줄
_KEY_LINE = re.compile(r"^[-*+•·]?\s*\*{0,2}([^*:：]{1,24})\*{0,2}\s*[:：]\s*(.*)$")


def _block_to_bullets(block: str) -> Dict:
    """개요 블록 하나 → bullets 슬라이드 dict(아트디렉터 실패 시 폴백).

    실제 개요는 이렇게 생겼다 — '본문 개요' 아래에 들여쓴 하위 불릿이 오고,
    그 뒤에 설계 메타 줄이 더 붙는다:

        - **레이아웃 제안**: 콘텐츠
        - **핵심 메시지(1개)**: …
        - **본문 개요**:
            - 정의: AECT …
            - '학습 촉진'과 '성과 개선'이 최종 목표
        - **시각자료 제안**: …
        - **발표자 노트(대화체)**: "자, 그럼 …"

    그래서 '본문 개요' 이후의 하위 불릿만 본문으로 삼고, 다음 라벨 줄이 나오면 멈춘다.
    """
    raw = [l for l in block.splitlines() if l.strip()]
    title = re.sub(r"^#{2,3}\s*", "", raw[0].strip()) if raw else "슬라이드"
    title = re.sub(r"^슬라이드\s*\d+\s*[—\-:：]\s*", "", title).strip()

    chip, level, bullets = "", "", []
    in_body = False
    for ln in raw[1:]:
        s = ln.strip()
        m = _KEY_LINE.match(s)
        key = m.group(1).strip() if m else ""
        val = m.group(2).strip() if m else ""

        if key.startswith("본문"):
            in_body = True
            if val:                     # 같은 줄에 인라인으로 쓴 경우
                bullets += [x.strip() for x in re.split(r"\s*[/·]\s*", val) if x.strip()]
            continue
        if m and any(k in key for k in _META_KEYS):
            if "지지" in key:           # '지지 목표: (분석) 목표2' → level 만 회수
                level = next((b for b in BLOOM_LEVELS if b in val), "")
            in_body = False             # 메타 줄이 나오면 본문 구간이 끝난 것
            continue
        if m and "핵심" in key:
            chip = val
            in_body = False
            continue

        # 라벨 없는 줄: 본문 구간 안이면 불릿, 밖이면 버린다
        if in_body:
            v = re.sub(r"^[-*+•·]\s*", "", s)
            if v:
                bullets.append(v)

    # 본문 구간을 못 찾았으면(형식이 다르면) 라벨 없는 줄을 본문으로 본다
    if not bullets:
        for ln in raw[1:]:
            s = ln.strip()
            m = _KEY_LINE.match(s)
            if m and (any(k in m.group(1) for k in _META_KEYS) or "핵심" in m.group(1)):
                continue
            v = re.sub(r"^[-*+•·]\s*", "", s)
            if v and not v.startswith("**"):
                bullets.append(v)

    return {"type": "bullets", "title": title, "chip": chip,
            "bullets": bullets[:5], "level": level}


def _fallback_plan(outline_md: str, deck_title: str) -> List[Dict]:
    """개요 md 를 단순 파싱해 bullets 위주 플랜으로(아트디렉터 실패 시)."""
    plan = [{"type": "cover", "title": deck_title, "chip": ""}]
    plan += [_block_to_bullets(b) for b in _outline_blocks(outline_md)]
    return plan


BLOOM_LEVELS = ("기억", "이해", "적용", "분석", "평가", "창조")


def _normalize_slide(s: Dict) -> Optional[Dict]:
    if not isinstance(s, dict):
        return None
    s.setdefault("type", "bullets")
    s.setdefault("title", "")
    s.setdefault("bullets", [])
    if s["type"] not in _RENDER:
        s["type"] = "bullets"
    if s["type"] == "cover":   # LLM이 표지를 만들면 내용형으로 강등
        s["type"] = "bullets"
    # 인지수준은 6수준 중 하나이거나 없음. 대괄호·설명이 붙어 와도 살려낸다.
    lv = str(s.get("level") or "")
    s["level"] = next((b for b in BLOOM_LEVELS if b in lv), "")
    return s


def level_counts(plan: List[Dict]) -> Dict[str, int]:
    """플랜의 인지수준 분포 — 교수자가 목표–자료 정렬을 눈으로 확인하는 값."""
    out: Dict[str, int] = {}
    for s in plan or []:
        lv = (s or {}).get("level") if isinstance(s, dict) else None
        if lv:
            out[lv] = out.get(lv, 0) + 1
    return {b: out[b] for b in BLOOM_LEVELS if b in out}


_CHUNK = 14  # 청크당 슬라이드 수(토큰 잘림 방지)


def plan_from_outline(generate_fn: Callable[[str, str, int], str],
                      outline_md: str, deck_title: str,
                      subtitle: str = "", course: str = "") -> List[Dict]:
    """개요를 슬라이드플랜으로 변환. 개요를 청크로 나눠 변환해 개수 잘림을 막는다.

    표지(cover)는 LLM이 아니라 여기서 맨 앞에 자동 추가한다. 청크가 실패하면
    그 청크만 블록 파싱(bullets)으로 폴백 → 전체 슬라이드 수는 항상 보존.
    """
    blocks = _outline_blocks(outline_md)
    if not blocks:
        plan = _fallback_plan(outline_md, deck_title)
        plan[0]["chip"] = subtitle
        plan[0]["band"] = course
        return plan

    chunks = [blocks[i:i + _CHUNK] for i in range(0, len(blocks), _CHUNK)]
    out: List[Dict] = []
    for ci, ch in enumerate(chunks):
        chunk_md = "\n\n".join(ch)
        user = (f"{_ART_RULES}\n\n[덱 제목]\n{deck_title}\n"
                f"[부분 {ci + 1}/{len(chunks)} · 이 부분의 슬라이드 {len(ch)}개를 빠짐없이 변환]\n\n"
                f"[슬라이드 개요]\n{chunk_md}")
        try:
            raw = generate_fn(_ART_SYS, user, 12000)
        except Exception as e:  # noqa: BLE001
            _log(f"[art] 청크 {ci + 1} 생성 오류: {e}")
            raw = ""
        arr = _extract_json_array(raw)
        got = [x for x in (_normalize_slide(s) for s in arr) if x] if isinstance(arr, list) else []
        # 개수가 모자라면 부족분을 블록 파싱으로 보충(개수 보존)
        if len(got) < len(ch):
            _log(f"[art] 청크 {ci + 1}: {len(got)}/{len(ch)} → 부족분 블록 폴백")
            got += [_block_to_bullets(b) for b in ch[len(got):]]
        out += got[:len(ch)]

    out.insert(0, {"type": "cover", "title": deck_title, "chip": subtitle,
                   "band": course})
    return out


def image_queries(plan: List[Dict]) -> Dict[int, str]:
    """photo 타입 슬라이드의 {인덱스: 영어 검색어}."""
    out = {}
    for i, s in enumerate(plan):
        if (s.get("type") == "photo") and s.get("image_query"):
            out[i] = s["image_query"]
    return out


# ── 이미지 생성 프롬프트 번들 (codex-prompt-img-studio 인풋 JSON) ──────────
# ★ 비율은 **정사각에 가깝게** 지시한다. 사진은 화면 전체가 아니라 좌우 세로 패널에
#   들어가고(4.18~4.78 × 4.69 in → 0.89~1.02), add_photo() 가 센터 크롭을 한다.
#   16:9 로 만들면 좌우가 32~41% 잘려 좌우 대조 구도가 통째로 날아간다(실측).
DEFAULT_STYLE_HINT = (
    "clean modern educational illustration, flat vector with subtle depth, "
    "consistent line weight, no gradient mesh, soft navy (#1E2761) and "
    "amber (#F2A900) accents, generous negative space, uncluttered, "
    "square 1:1 composition with the subject centered and margins on all sides "
    "(the image is center-cropped into a near-square panel)")

_IMGPROMPT_SYS = ("너는 강의 슬라이드용 이미지 생성 프롬프트 작가다. 각 슬라이드 제목/요약을 보고 "
                  "이미지 생성 모델(diffusion)용 **영어** 프롬프트 한 줄을 만든다. 그림 안에 글자·워터마크·"
                  "로고가 들어가지 않게 하고, 주제를 상징하는 구체적 장면/오브젝트로 묘사한다. "
                  "구도는 **정사각(1:1)에 가깝게** 잡고 주요 오브젝트를 가운데 모은다 — "
                  "좌우로 넓게 벌린 구도는 슬라이드에서 양옆이 잘린다. "
                  "**공통 스타일 문구는 반복하지 마라**(시스템이 뒤에 한 번 붙인다). "
                  "JSON 배열 [{\"n\":정수,\"prompt\":\"...\"}] 만 출력한다.")

# 이미지를 실제로 '올릴 자리'가 있는 타입.
#   photo   — 원래부터 사진 슬라이드
#   bullets — 이미지가 들어오면 photo 로 승격(칩·불릿이 그대로 옆으로 간다)
# cards/process/compare/table/quiz/objectives/agenda/closing/stat/section/cover 는
# 도형 레이아웃이 자리를 다 쓰므로 이미지를 넣지 않는다.
PLACEABLE_TYPES = ("photo", "bullets")


def _prompt_place(s: Dict) -> bool:
    """이 슬라이드에 이미지를 실제로 배치할 것인가(프롬프트 JSON 의 place 값).

    ★ 예전에는 emphasis 슬라이드를 빼고 있었다. 그런데 emphasis 는 **칩 색만**
      바꾼다(연한 칩 → 진한 네이비 알약). 이미지 자리와 아무 상관이 없고,
      render_photo 는 emphasis 여도 사진 자리를 그대로 예약한다.
      그래서 '자리는 예약됐는데 프롬프트가 없는' 슬라이드가 생겼고,
      받는 사람은 만들 수 없는 이미지를 기다리다 빈 액자로 인쇄했다.
      (3주차 7번·76번이 그랬다)

    지켜야 할 불변식: **자리를 예약하는 슬라이드에는 반드시 프롬프트가 있다.**
    자리는 type=photo 일 때 예약되고, bullets 는 이미지가 오면 photo 로 승격되므로
    둘 다 프롬프트 대상이다 = PLACEABLE_TYPES.
    """
    return (s.get("type") or "bullets") in PLACEABLE_TYPES


def photo_slot_slides(plan: List[Dict]) -> List[int]:
    """사진 **자리가 예약되는** 슬라이드의 1-based 번호.

    빈 액자로 남을 수 있는 자리를 세는 데 쓴다 — 조용히 비면 인쇄까지 간다.
    """
    return [i + 1 for i, s in enumerate(plan or [])
            if (s.get("type") or "bullets") == "photo"]


def apply_images(plan: List[Dict], images: Dict[int, bytes]):
    """외부 이미지를 플랜에 붙인다. bullets 슬라이드는 photo 로 승격.

    반환: (새 플랜, 사용된 {인덱스: 바이트}, 자리 없어 건너뛴 1-based 번호 목록)
    이미지를 조용히 버리지 않고 무엇이 안 들어갔는지 항상 돌려준다.
    """
    out, used, skipped = [], {}, []
    for i, s in enumerate(plan or []):
        s = dict(s) if isinstance(s, dict) else {"type": "bullets", "title": str(s)}
        if i in (images or {}):
            typ = (s.get("type") or "bullets").lower()
            if typ in PLACEABLE_TYPES:
                s["type"] = "photo"
                used[i] = images[i]
            else:
                skipped.append(i + 1)
        out.append(s)
    return out, used, skipped


def image_prompt_bundle(plan: List[Dict], deck_title: str, *,
                        generate_fn: Optional[Callable[[str, str, int], str]] = None,
                        style_hint: str = DEFAULT_STYLE_HINT,
                        have_photos: Optional[set] = None) -> Dict:
    """**이미지를 실제로 올릴 수 있는 슬라이드만** 프롬프트로 뽑아 1개 번들(dict)로.

    전 슬라이드를 뽑던 것을 좁힌 이유: 배치 가능한 타입은 photo·bullets 뿐이고
    (나머지는 도형 레이아웃이 자리를 다 쓴다), 42장을 뽑으면 34장이 '자리 없어
    건너뜀' 으로 버려졌다. 만들지도 못할 프롬프트에 LLM 호출을 3배 쓰고,
    받는 사람은 42장을 그리고 34장을 헛수고했다.

    generate_fn 이 있으면 LLM으로 영어 프롬프트를 짓고(청크·폴백), 없으면
    image_query/제목 기반으로 결정론적으로 만든다. 'prompts' 배열은 위→아래
    순차 실행용이고, "n" 은 **1-based 슬라이드 번호** = 이미지 파일명 규칙의
    번호다(003.png → 3번 슬라이드). 그래서 걸러내도 번호는 원래 것을 유지한다.
    """
    total = len(plan)
    # 자동 사진이 이미 붙은 슬롯은 뺀다 — 그건 그릴 필요가 없다.
    # 남는 것 = "자리는 있는데 사진을 못 찾은 슬라이드" 뿐이고, 그게 그려야 할 목록이다.
    got = set(int(i) for i in (have_photos or ()))
    idx = [i for i, s_ in enumerate(plan) if _prompt_place(s_) and i not in got]
    n = len(idx)
    en: Dict[int, str] = {}
    if generate_fn and idx:
        for a in range(0, n, _CHUNK):
            ch = idx[a:a + _CHUNK]
            lines = "\n".join(
                f'{i + 1}. {plan[i].get("title", "")} | {plan[i].get("chip", "")}' for i in ch)
            user = (f"[공통 스타일] {style_hint}\n"
                    f"아래 슬라이드 {len(ch)}개 각각에 대해 n(슬라이드 번호)과 영어 이미지 프롬프트를 만들어라. "
                    f"프롬프트에 공통 스타일을 녹이고, 그림 안 텍스트 금지.\n\n{lines}")
            try:
                arr = _extract_json_array(generate_fn(_IMGPROMPT_SYS, user, 6000))
            except Exception as e:  # noqa: BLE001
                _log(f"[imgprompt] 청크 오류: {e}")
                arr = None
            if isinstance(arr, list):
                for o in arr:
                    if isinstance(o, dict) and o.get("n") and o.get("prompt"):
                        en[int(o["n"]) - 1] = str(o["prompt"]).strip()

    prompts = []
    for i in idx:
        s = plan[i]
        subj = (en.get(i) or s.get("image_query")
                or f"conceptual illustration for '{s.get('title', '')}'")
        subj = str(subj).strip().rstrip(".").strip()       # 마침표 중복 방지
        # LLM 이 공통 스타일을 이미 녹여 넣는 경우가 있다 — 그대로 또 붙이면
        # 같은 문구가 두 번 들어가 프롬프트가 두 배로 길어지고 주제가 묻힌다.
        tail = "" if "flat vector" in subj.lower() else f" {style_hint}."
        prompt = f"{subj}.{tail} No text, no watermark, no logo."
        prompts.append({
            "n": i + 1,
            "title": s.get("title", ""),
            "type": s.get("type", "bullets"),
            # 인지수준을 함께 흘린다 — 이 번들이 슬라이드→비디오 단계의 입력이라,
            # 씬별 나레이션 길이·속도를 정할 근거가 된다(기억<이해<분석 순으로 체류).
            "level": s.get("level", ""),
            "prompt": prompt,
            "negative": "text, letters, watermark, logo, low quality, distorted",
            "keywords": ([s["image_query"]] if s.get("image_query") else []),
            "place": _prompt_place(s),
        })
    return {
        "deck": deck_title,
        "style_hint": style_hint,
        # 슬라이드의 사진 자리는 세로 패널(≈1:1)이다. landscape 로 두면
        # 생성 도구가 3:2 를 골라 좌우 30~40% 가 잘린다.
        "aspect": "square",
        "target_box": "near-square vertical panel (~1:1), center-cropped",
        "count": n,                 # 뽑은 프롬프트 수 = 그려야 할 장수
        "deck_slides": total,       # 덱 전체 장수 (번호 n 은 이 기준이다)
        "photos_found": len(got),   # 자동 사진으로 이미 채워진 장수(프롬프트에서 제외됨)
        "file_naming": ("생성한 이미지는 images/ 폴더에 '슬라이드번호'로 저장하세요. "
                        "예: 003.png → 3번 슬라이드. 파일명 앞 숫자만 맞으면 되고 "
                        "확장자·뒤 설명은 자유입니다(003_뇌구조.png)."),
        "prompts": prompts,
    }
