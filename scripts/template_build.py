# -*- coding: utf-8 -*-
r"""레이아웃 견본 PPTX 여러 개 → templates/*.json 한 벌.

    .venv-app\Scripts\python scripts\template_build.py _context\6 --wipe

견본 6종은 **골격이 같고 색만 다르다.** 그래서 색을 손으로 고르지 않고
'어느 자리에 쓰였나'로 규칙을 세워 뽑는다 — 견본이 늘어도 같은 규칙이 돈다.

자리 판별은 **면적**으로 한다. 6종이 전부 아래 서명을 갖는다:

    106.9 in²  중립 워시     (F1F5F9 — 페이지·존 바탕)
     72.6 in²  흰색
     45.0 in²  ★ 주색        (표지 밴드·섹션 — deck 의 navy 자리)
     21.8 in²  ★ 주색 워시   (요약 존 바탕 — deck 의 chip 자리)
     15.5 in²  ★ 보조색      (작은 강조 30개 — deck 의 amber 자리)
      5.2 in²  ★ 중간 톤     (deck 의 navy_55 자리)

글자색 중 가장 어두운 주색 계열이 밴드용 진한 색(navy_dk)이다. 없으면 만든다.

★ 대비는 만들고 나서 **검사한다.** 색을 바꾸면 흰 글씨가 안 읽히는 밴드가
  반드시 생긴다(앰버 #F2A900 에 흰 글씨 2.1:1 이 그랬다). 기준에 못 미치면
  진한 짝을 자동으로 어둡게 해서 통과시킨다.
"""
from __future__ import annotations

import io
import json
import re
import sys
from collections import Counter, defaultdict
from pathlib import Path

out = io.open(sys.stdout.fileno(), "w", encoding="utf-8", closefd=False)
ROOT = Path(__file__).resolve().parents[1]
TPL_DIR = ROOT / "templates"

# 면적 서명 → 역할. 견본마다 소수점이 흔들리므로 가까운 것을 고른다.
SLOTS = [(45.0, "navy"), (21.8, "chip"), (15.5, "amber"), (5.2, "navy_55")]
NEUTRAL = {"F1F5F9", "FFFFFF", "E2E8F0", "CBD5E1", "475569", "94A3B8", "334155"}

# 견본이 공통으로 쓰는 중립색 — 색상환과 무관하므로 그대로 가져간다.
FIXED = {"navy_30": "CBD5E1", "navy_12": "E2E8F0",
         "on_navy": "F1F5F9", "white": "FFFFFF", "ink": "334155"}

LABEL = {
    "cyan": "시안", "green": "그린", "khaki": "카키",
    "lavender": "라벤더", "orange": "오렌지", "wine": "와인",
}


# ── 색 계산 ────────────────────────────────────────────────────────────────
def rgb(h: str) -> tuple:
    h = h.lstrip("#")
    return tuple(int(h[i:i + 2], 16) for i in (0, 2, 4))


def hexs(t) -> str:
    return "".join(f"{max(0, min(255, round(v))):02X}" for v in t)


def lum(h: str) -> float:
    def ch(v):
        v /= 255.0
        return v / 12.92 if v <= 0.04045 else ((v + 0.055) / 1.055) ** 2.4
    r, g, b = (ch(v) for v in rgb(h))
    return 0.2126 * r + 0.7152 * g + 0.0722 * b


def ratio(a: str, b: str) -> float:
    la, lb = lum(a), lum(b)
    hi, lo = max(la, lb), min(la, lb)
    return (hi + 0.05) / (lo + 0.05)


def mix(a: str, b: str, t: float) -> str:
    """a 에서 b 쪽으로 t 만큼. t=0 → a, t=1 → b."""
    return hexs(x + (y - x) * t for x, y in zip(rgb(a), rgb(b)))


def darken_until(h: str, against: str, want: float) -> str:
    """대비가 want 를 넘을 때까지 검정 쪽으로 섞는다. 색상은 유지된다."""
    for i in range(0, 96, 2):
        c = mix(h, "000000", i / 100.0)
        if ratio(c, against) >= want:
            return c
    return "000000"


# ── 추출 ───────────────────────────────────────────────────────────────────
def _srgb(el):
    from pptx.oxml.ns import qn
    for c in el.iter(qn("a:srgbClr")):
        v = c.get("val")
        if v and re.fullmatch(r"[0-9A-Fa-f]{6}", v):
            return v.upper()
    return None


def scan(path: Path) -> dict:
    from pptx import Presentation
    from pptx.oxml.ns import qn
    prs = Presentation(str(path))
    area = defaultdict(float)
    text = Counter()
    fonts = Counter()
    for s in prs.slides:
        for sh in s.shapes:
            spPr = sh._element.find(qn("p:spPr"))
            if spPr is not None:
                solid = spPr.find(qn("a:solidFill"))
                if solid is not None and (v := _srgb(solid)):
                    area[v] += ((sh.width or 0) / 914400) * ((sh.height or 0) / 914400)
            if sh.has_text_frame:
                for p in sh.text_frame.paragraphs:
                    for r in p.runs:
                        rPr = r._r.find(qn("a:rPr"))
                        if rPr is None:
                            continue
                        if (v := _srgb(rPr)):
                            text[v] += 1
                        latin = rPr.find(qn("a:latin"))
                        if latin is not None and latin.get("typeface"):
                            fonts[latin.get("typeface")] += 1
    return {"area": dict(area), "text": text, "font": (fonts.most_common(1) or [("", 0)])[0][0]}


def pick_slots(area: dict) -> dict:
    """면적 서명으로 자리를 정한다. 중립색은 후보에서 뺀다."""
    cand = {c: a for c, a in area.items() if c not in NEUTRAL}
    got, used = {}, set()
    for want, role in SLOTS:
        best, gap = None, 1e9
        for c, a in cand.items():
            if c in used:
                continue
            if (d := abs(a - want)) < gap:
                best, gap = c, d
        if best is None:
            raise SystemExit(f"자리 '{role}' 를 못 찾았습니다. 견본 구조가 다릅니다.")
        got[role], _ = best, used.add(best)
    return got


def build(path: Path) -> dict:
    d = scan(path)
    slot = pick_slots(d["area"])
    navy, amber, chip = slot["navy"], slot["amber"], slot["chip"]

    # 밴드에 흰 글씨를 얹으려면 4.5:1 이 필요하다. 견본 글자색 중 주색 계열
    # 가장 어두운 것을 먼저 보고, 그래도 모자라면 어둡게 만든다.
    dark_cands = [c for c in d["text"] if c not in NEUTRAL and lum(c) < lum(navy)]
    navy_dk = min(dark_cands, key=lum) if dark_cands else navy
    if ratio(navy_dk, "FFFFFF") < 4.5:
        navy_dk = darken_until(navy, "FFFFFF", 4.5)

    colors = {
        "navy": navy,
        "navy_75": mix(navy, "FFFFFF", 0.25),
        "navy_55": slot["navy_55"],
        "chip": chip,
        "amber": amber,
        "amber_dk": darken_until(amber, "FFFFFF", 4.5),
        "amber_wash": mix(amber, "FFFFFF", 0.88),
        "navy_dk": navy_dk,
        **FIXED,
    }
    key = path.stem.split("_")[-1].lower()
    return {
        "name": key,
        "label": f"{LABEL.get(key, key)} ({colors['navy'].lower()})",
        "note": f"{path.name} 에서 뽑은 배색. 6종은 골격이 같고 색만 다르다 — "
                f"면적 서명(45.0/21.8/15.5/5.2 in²)으로 자리를 정하고, "
                f"밴드용 진한 색은 흰 글씨 4.5:1 을 넘도록 계산했다.",
        "font": d["font"] or "Noto Sans KR",
        "source": f"_context/6/{path.name}",
        "colors": colors,
        "_mapping": {
            "navy": "표지 밴드·섹션·제목 (면적 45.0 자리)",
            "chip": "요약 존 배경 (면적 21.8 자리)",
            "amber": "작은 강조 — 룰·점 (면적 15.5 자리, 30개)",
            "navy_55": "보조 글자·중간 톤 (면적 5.2 자리)",
            "navy_dk": "흰 글씨를 얹는 밴드 배경 — 주색이 4.5:1 에 못 미칠 때 쓴다",
            "amber_dk": "보조색 밴드 배경 — 같은 이유",
            "ink": "본문 글자 (slate-700)",
        },
    }


def main() -> int:
    args = [a for a in sys.argv[1:] if not a.startswith("--")]
    if not args:
        out.write(__doc__ + "\n")
        return 2
    src = Path(args[0]).resolve()
    files = sorted(src.glob("*.pptx")) if src.is_dir() else [src]
    if not files:
        out.write(f"pptx 가 없습니다: {src}\n")
        return 2

    TPL_DIR.mkdir(parents=True, exist_ok=True)
    if "--wipe" in sys.argv:
        old = sorted(TPL_DIR.glob("*.json"))
        for f in old:
            f.unlink()
        out.write(f"기존 템플릿 {len(old)}개 삭제: "
                  f"{', '.join(f.stem for f in old) or '없음'}\n\n")

    for f in files:
        t = build(f)
        (TPL_DIR / f"{t['name']}.json").write_text(
            json.dumps(t, ensure_ascii=False, indent=2), encoding="utf-8")
        c = t["colors"]
        out.write(f"{t['name']:10s} {t['label']:18s} {t['font']}\n")
        out.write(f"   navy #{c['navy']}  amber #{c['amber']}  chip #{c['chip']}\n")
        out.write(f"   밴드 navy_dk #{c['navy_dk']} (흰글씨 {ratio(c['navy_dk'], 'FFFFFF'):.1f}:1)"
                  f"  amber_dk #{c['amber_dk']} ({ratio(c['amber_dk'], 'FFFFFF'):.1f}:1)\n")
    out.write(f"\n-> {TPL_DIR} 에 {len(files)}개\n")
    return 0


if __name__ == "__main__":
    sys.exit(main())
