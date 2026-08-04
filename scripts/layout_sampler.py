# -*- coding: utf-8 -*-
r"""레이아웃 견본 뽑기 — 실제로 만들어진 덱에서 **종류별 한 장씩**만 남긴다.

    .venv\Scripts\python scripts\layout_sampler.py <덱.pptx> <슬라이드플랜.json> [출력.pptx]

디자인을 손보려면 13종이 한 파일에 모여 있어야 비교가 된다. 코드로 새로 그리지
않고 **실물 덱에서 골라내는** 이유: 실제 좌표·폰트·이미지가 그대로 들어 있어야
"이 사진 높이가 맞나" 같은 판단을 할 수 있다.

python-pptx 는 슬라이드 삭제 API 가 없어서 sldIdLst 에서 직접 뺀다.
"""
from __future__ import annotations

import json
import sys
from collections import OrderedDict
from pathlib import Path


def main() -> int:
    if len(sys.argv) < 3:
        print(__doc__)
        return 2
    deck = Path(sys.argv[1]).resolve()
    plan_p = Path(sys.argv[2]).resolve()
    out = Path(sys.argv[3]).resolve() if len(sys.argv) > 3 else \
        deck.parent / "레이아웃-견본.pptx"

    from pptx import Presentation
    from pptx.oxml.ns import qn

    plan = json.loads(plan_p.read_text(encoding="utf-8"))
    if isinstance(plan, dict):
        plan = plan.get("slides") or []

    prs = Presentation(str(deck))
    n_deck = len(prs.slides._sldIdLst)
    if len(plan) != n_deck:
        print(f"[경고] 플랜 {len(plan)}장 != 덱 {n_deck}장 — 앞쪽 기준으로 맞춥니다.")

    # 종류별 **첫 등장** 한 장만. 사진이 실제로 붙은 것을 우선 고른다.
    pick: "OrderedDict[str, int]" = OrderedDict()
    for i, s in enumerate(plan[:n_deck]):
        t = (s.get("type") or "bullets") if isinstance(s, dict) else "bullets"
        if t not in pick:
            pick[t] = i
        elif t == "photo" and s.get("_credit_short"):
            pick[t] = i          # 사진이 붙은 쪽으로 교체

    keep = sorted(pick.values())
    print(f"레이아웃 {len(pick)}종 · 남길 슬라이드 {[i + 1 for i in keep]}")
    for t, i in pick.items():
        title = (plan[i].get("title") or "")[:34]
        print(f"  {i + 1:3d}  {t:11s} {title}")

    # ★ 뒤에서부터 지운다 — 앞에서 지우면 인덱스가 밀려 엉뚱한 장이 빠진다.
    part = prs.part
    ids = prs.slides._sldIdLst
    for i in range(n_deck - 1, -1, -1):
        if i in keep:
            continue
        sid = list(ids)[i]
        rid = sid.get(qn("r:id"))
        if rid and rid in part.rels:
            try:
                part.drop_rel(rid)
            except Exception:      # noqa: BLE001 — 이미 정리된 관계는 넘어간다
                pass
        ids.remove(sid)

    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    print(f"\n-> {out}  ({out.stat().st_size / 1048576:.1f} MB · {len(keep)}장)")
    return 0


if __name__ == "__main__":
    sys.exit(main())
