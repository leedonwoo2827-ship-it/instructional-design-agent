# -*- coding: utf-8 -*-
r"""템플릿 PPTX → 색·폰트 토큰 추출.

    .venv-app\Scripts\python scripts\template_extract.py <템플릿.pptx> [--save 이름]

배자(좌표)는 읽지 않는다. 내용 길이가 템플릿과 다를 때 넘치거나 잉여 공간이
생기기 때문이고, 그건 코드의 _fit() 이 이미 잘 처리한다. 여기서 가져오는 것은
**색과 글꼴** 뿐이다 — 그것만 바뀌어도 덱의 인상이 완전히 달라진다.

쓰인 횟수와 면적으로 정렬해서 보여준다. 큰 도형에 많이 쓰인 색이 그 템플릿의
주색이다.
"""
from __future__ import annotations

import io
import json
import re
import sys
from collections import Counter, defaultdict
from pathlib import Path

out = io.open(sys.stdout.fileno(), "w", encoding="utf-8", closefd=False)
ROOT = Path(__file__).resolve().parents[1]
TPL_DIR = ROOT / "templates"


def _srgb(el) -> str | None:
    """도형/글자 채움에서 srgbClr 값을 캔다. 테마색 참조는 건너뛴다."""
    from pptx.oxml.ns import qn
    for c in el.iter(qn("a:srgbClr")):
        v = c.get("val")
        if v and re.fullmatch(r"[0-9A-Fa-f]{6}", v):
            return v.upper()
    return None


def extract(path: Path) -> dict:
    from pptx import Presentation
    from pptx.oxml.ns import qn

    prs = Presentation(str(path))
    fill_area: dict[str, float] = defaultdict(float)   # 색 → 총 면적(in²)
    fill_n: Counter = Counter()
    text_n: Counter = Counter()
    fonts: Counter = Counter()
    sizes: Counter = Counter()

    for s in prs.slides:
        for sh in s.shapes:
            # 도형 채움
            spPr = sh._element.find(qn("p:spPr"))
            if spPr is not None:
                solid = spPr.find(qn("a:solidFill"))
                if solid is not None and (v := _srgb(solid)):
                    w = (sh.width or 0) / 914400
                    h = (sh.height or 0) / 914400
                    fill_area[v] += w * h
                    fill_n[v] += 1
            # 글자 색·글꼴·크기
            if sh.has_text_frame:
                for p in sh.text_frame.paragraphs:
                    for r in p.runs:
                        rPr = r._r.find(qn("a:rPr"))
                        if rPr is not None:
                            if (v := _srgb(rPr)):
                                text_n[v] += 1
                            latin = rPr.find(qn("a:latin"))
                            if latin is not None and latin.get("typeface"):
                                fonts[latin.get("typeface")] += 1
                            if rPr.get("sz"):
                                sizes[int(rPr.get("sz")) // 100] += 1

    return {
        "source": path.name,
        "slide_size_in": [round(prs.slide_width / 914400, 2),
                          round(prs.slide_height / 914400, 2)],
        "fills_by_area": sorted(((c, round(a, 1), fill_n[c]) for c, a in fill_area.items()),
                                key=lambda t: -t[1]),
        "text_colors": text_n.most_common(),
        "fonts": fonts.most_common(),
        "sizes_pt": sorted(sizes.items(), key=lambda t: -t[1]),
    }


def main() -> int:
    if len(sys.argv) < 2:
        out.write(__doc__ + "\n")
        return 2
    path = Path(sys.argv[1]).resolve()
    d = extract(path)

    out.write(f"== {d['source']} ==\n")
    out.write(f"슬라이드 크기 {d['slide_size_in'][0]} x {d['slide_size_in'][1]} in\n")

    out.write("\n[도형 채움색 — 면적 큰 순]\n")
    for c, area, n in d["fills_by_area"][:14]:
        out.write(f"  #{c}  {area:7.1f} in²  {n:3d}개\n")

    out.write("\n[글자색 — 많이 쓰인 순]\n")
    for c, n in d["text_colors"][:10]:
        out.write(f"  #{c}  {n:4d}회\n")

    out.write("\n[글꼴]\n")
    for f, n in d["fonts"][:8]:
        out.write(f"  {f:24s} {n:4d}회\n")

    out.write("\n[글자 크기(pt)]\n")
    out.write("  " + " · ".join(f"{s}pt×{n}" for s, n in d["sizes_pt"][:10]) + "\n")

    if "--save" in sys.argv:
        i = sys.argv.index("--save")
        name = sys.argv[i + 1] if len(sys.argv) > i + 1 else path.stem
        TPL_DIR.mkdir(parents=True, exist_ok=True)
        p = TPL_DIR / f"{name}.raw.json"
        p.write_text(json.dumps(d, ensure_ascii=False, indent=2), encoding="utf-8")
        out.write(f"\n-> {p}\n")
    return 0


if __name__ == "__main__":
    sys.exit(main())
